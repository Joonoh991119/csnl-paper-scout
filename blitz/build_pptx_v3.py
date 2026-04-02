"""
Paper Blitz PPT Builder v3 — clean academic design, figure-first, aspect-ratio preserved.
"""

import os
from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Design tokens ───────────────────────────────────────────────

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xFA, 0xFA, 0xFA)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
DARK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xAA, 0xAA, 0xAA)
ACCENT = RGBColor(0x00, 0x72, 0xB2)  # Nature blue

FONT = "Helvetica Neue"

SW = Inches(13.333)
SH = Inches(7.5)


# ── Helpers ─────────────────────────────────────────────────────

def _text(slide, text, left, top, width, height,
          size=16, color=BLACK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    for run in p.runs:
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
    return box


def _multi(slide, lines, left, top, width, height,
           size=14, color=BLACK, bold=False, line_spacing=1.3, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, (line, line_bold, line_color) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = Pt(size * line_spacing)
        if not line.strip():
            p.space_before = Pt(4)
            continue
        p.text = line
        for run in p.runs:
            run.font.name = FONT
            run.font.size = Pt(size)
            run.font.color.rgb = line_color if line_color else color
            run.font.bold = line_bold if line_bold is not None else bold
    return box


def _img_ar(slide, path, left, top, max_w, max_h):
    """Add image preserving aspect ratio within max bounds."""
    if not os.path.exists(path):
        print(f"  MISSING: {path}")
        return None
    im = Image.open(path)
    w_px, h_px = im.size
    ratio = w_px / h_px

    # Fit within max_w × max_h preserving aspect ratio
    w = max_w
    h = w / ratio
    if h > max_h:
        h = max_h
        w = h * ratio

    return slide.shapes.add_picture(path, left, top, int(w), int(h))


def _bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


# ── Slides ──────────────────────────────────────────────────────

def slide_1_title(prs, fig_dir):
    """Title slide: header crop + why this paper."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)

    # Title header image (top, wide, aspect preserved)
    _img_ar(slide, f"{fig_dir}/title_header.png",
            Inches(0.5), Inches(0.4), Inches(12.3), Inches(2.2))

    # Thin separator
    sep = slide.shapes.add_shape(1, Inches(0.5), Inches(2.8), Inches(12.3), Pt(0.75))
    sep.fill.solid()
    sep.fill.fore_color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    sep.line.fill.background()

    # Why this paper section
    _text(slide, "Why this paper?",
          Inches(0.5), Inches(3.2), Inches(4), Inches(0.5),
          size=16, color=ACCENT, bold=True)

    _multi(slide, [
        ("JOP — asymmetric prior × time estimation", True, DARK),
        ("Prior distribution shape → relative-to-absolute scale transformation (normative model)", False, GRAY),
        ("", None, None),
        ("This paper shows that prior shape (bimodal) produces qualitatively distinct", False, DARK),
        ("estimation patterns in naturalistic sensorimotor tasks — not just linear bias.", False, DARK),
        ("Prior reliance scales with uncertainty, consistent with Bayesian integration.", False, DARK),
        ("", None, None),
        ("→ Empirical support for shape-dependent, nonlinear estimation effects", True, ACCENT),
        ("   that JOP's normative model predicts in the temporal domain.", False, ACCENT),
    ], Inches(0.5), Inches(3.8), Inches(12), Inches(3.5),
       size=13, color=DARK, line_spacing=1.4)


def slide_2_background(prs, fig_dir):
    """Background & Hypothesis: Fig 1 (Bayesian simulation) + gap + predictions."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)

    _text(slide, "Background & Hypothesis",
          Inches(0.5), Inches(0.25), Inches(8), Inches(0.5),
          size=20, bold=True)

    # Left column: background text
    _multi(slide, [
        ("Established", True, DARK),
        ("Bayesian integration: sensory + prior → posterior, weighted by reliability", False, GRAY),
        ("Robust in controlled lab tasks (reaching, pointing, force estimation)", False, GRAY),
        ("", None, None),
        ("Gap", True, DARK),
        ("Untested in naturalistic tasks with non-Gaussian priors", False, GRAY),
        ("and complex, full-body movements", False, GRAY),
        ("", None, None),
        ("Predictions", True, ACCENT),
        ("P1: If bimodal prior is learned → nonlinear jump in estimation error", False, DARK),
        ("     between left/right distribution segments", False, DARK),
        ("P2: If Bayesian weighting → jump magnitude scales with uncertainty", False, DARK),
    ], Inches(0.4), Inches(0.9), Inches(5.2), Inches(6.2),
       size=11, line_spacing=1.35)

    # Right: Fig 1 (full panels A-D)
    _img_ar(slide, f"{fig_dir}/fig1_full.png",
            Inches(5.8), Inches(0.5), Inches(7.2), Inches(6.8))

    # Figure annotation
    _multi(slide, [
        ("Fig. 1 — Bayesian simulation predictions", False, LIGHT_GRAY),
        ("Green = slow (low σ)  Blue = moderate  Red = fast (high σ)", False, LIGHT_GRAY),
        ("x: ball position (cm)  y: estimation error (cm)", False, LIGHT_GRAY),
    ], Inches(5.8), Inches(6.5), Inches(7), Inches(1),
       size=9, line_spacing=1.2)


def slide_3_paradigm(prs, fig_dir):
    """Experimental paradigm: generated diagram + XR setup photos."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)

    _text(slide, "Experimental Paradigm",
          Inches(0.5), Inches(0.25), Inches(8), Inches(0.5),
          size=20, bold=True)

    # XR setup photos (top right, compact)
    _img_ar(slide, f"{fig_dir}/fig2_xr_photos.png",
            Inches(7.5), Inches(0.15), Inches(5.7), Inches(2.5))

    _text(slide, "Fig. 2 — XR tennis serve return in life-sized CAVE",
          Inches(7.5), Inches(2.55), Inches(5.5), Inches(0.3),
          size=9, color=LIGHT_GRAY)

    # Paradigm diagram (center-bottom, large)
    _img_ar(slide, f"{fig_dir}/paradigm_diagram.png",
            Inches(0.3), Inches(2.9), Inches(12.7), Inches(4.4))

    # Key rationale text (top left, below title)
    _multi(slide, [
        ("Rationale: ball speed manipulates visual uncertainty", True, DARK),
        ("while temporal demand (400 ms window) stays constant", False, GRAY),
        ("→ isolates the effect of sensory reliability on prior integration", False, GRAY),
    ], Inches(0.5), Inches(0.85), Inches(6.8), Inches(1.8),
       size=11, line_spacing=1.3)


def slide_4_results(prs, fig_dir):
    """Key Results: Fig 3 large center + figure annotations + findings."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)

    _text(slide, "Key Results",
          Inches(0.5), Inches(0.25), Inches(6), Inches(0.5),
          size=20, bold=True)

    # Fig 3 — large, center
    _img_ar(slide, f"{fig_dir}/fig3_results.png",
            Inches(0.8), Inches(0.9), Inches(11.5), Inches(4.5))

    # Figure annotation (right of figure)
    _multi(slide, [
        ("Fig. 3 — Bimodal prior effect × visual uncertainty", False, LIGHT_GRAY),
        ("x: true ball position (cm)   y: estimation error (cm)", False, LIGHT_GRAY),
        ("Green = slow   Blue = moderate   Red = fast", False, LIGHT_GRAY),
        ("Black arrows = magnitude of bimodal jump between segments", False, LIGHT_GRAY),
    ], Inches(0.8), Inches(5.3), Inches(7), Inches(1),
       size=9, line_spacing=1.2)

    # Key findings (bottom, compact)
    _multi(slide, [
        ("After consolidation (Days 2-3)", True, DARK),
        ("Moderate: B = −2.64, p = .006     Fast: B = −2.36, p = .023     Slow: B = −0.81, p = .154 (n.s.)", False, GRAY),
        ("No effect on Day 1 → overnight consolidation required   |   No explicit awareness (implicit learning)", False, GRAY),
    ], Inches(0.5), Inches(6.0), Inches(12.3), Inches(1.3),
       size=11, line_spacing=1.35)


def slide_5_model(prs, fig_dir):
    """Model: Fig 4 + Fig 5 side by side with explanation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)

    _text(slide, "Combined Model: Bayesian Prior + Biomechanical Costs",
          Inches(0.5), Inches(0.25), Inches(12), Inches(0.5),
          size=20, bold=True)

    # Fig 4 left
    _img_ar(slide, f"{fig_dir}/fig4_biomech.png",
            Inches(0.3), Inches(1.0), Inches(6.2), Inches(3.5))

    # Fig 5 right
    _img_ar(slide, f"{fig_dir}/fig5_combined.png",
            Inches(6.8), Inches(1.0), Inches(6.2), Inches(3.5))

    # Labels
    _text(slide, "Fig. 4 — Control exp. (uniform distribution): biomechanical bias",
          Inches(0.3), Inches(4.5), Inches(6), Inches(0.3),
          size=9, color=LIGHT_GRAY)
    _text(slide, "Fig. 5 — Combined model: prior + biomechanics → qualitative match",
          Inches(6.8), Inches(4.5), Inches(6), Inches(0.3),
          size=9, color=LIGHT_GRAY)

    # Explanation
    _multi(slide, [
        ("Control experiment (uniform prior) reveals systematic linear bias from motor costs", False, DARK),
        ("Combined model = Bayesian prior integration + biomechanical bias term from control parameters", False, DARK),
        ("→ Qualitatively matches experimental data: behavior is jointly shaped by probabilistic inference and motor costs", False, ACCENT),
    ], Inches(0.5), Inches(5.1), Inches(12), Inches(2.2),
       size=12, line_spacing=1.4)

    # Axis annotations
    _multi(slide, [
        ("Both: x = ball position (cm), y = estimation error (cm)", False, LIGHT_GRAY),
        ("Green = slow  Blue = moderate  Red = fast", False, LIGHT_GRAY),
    ], Inches(0.3), Inches(6.5), Inches(6), Inches(0.8),
       size=9, line_spacing=1.2)


def slide_6_takeaway(prs, fig_dir):
    """Takeaway: 3 core messages + key figure."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)

    _text(slide, "Takeaway",
          Inches(0.5), Inches(0.25), Inches(4), Inches(0.5),
          size=22, bold=True)

    # Core messages (left)
    _multi(slide, [
        ("1. Implicit bimodal prior learning in naturalistic sensorimotor tasks", True, DARK),
        ("   Humans acquire complex (non-Gaussian) environmental statistics", False, GRAY),
        ("   without conscious awareness — even in full-body, ecologically valid tasks", False, GRAY),
        ("", None, None),
        ("2. Prior reliance scales with sensory uncertainty (Bayesian integration)", True, DARK),
        ("   Moderate & fast conditions show significant bimodal jump;", False, GRAY),
        ("   slow condition does not — consistent with reliability weighting", False, GRAY),
        ("", None, None),
        ("3. Behavior = Bayesian inference + biomechanical costs", True, DARK),
        ("   Combined model qualitatively captures the full data pattern", False, GRAY),
        ("   → prior shape + motor constraints jointly determine estimation", False, GRAY),
    ], Inches(0.4), Inches(1.0), Inches(6), Inches(6),
       size=12, line_spacing=1.25)

    # Key figure (right)
    _img_ar(slide, f"{fig_dir}/fig3_results.png",
            Inches(6.5), Inches(1.0), Inches(6.5), Inches(5.5))

    _text(slide, "Fig. 3",
          Inches(6.5), Inches(6.6), Inches(3), Inches(0.3),
          size=9, color=LIGHT_GRAY)


# ── Main ────────────────────────────────────────────────────────

def build(fig_dir, out_path):
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    slide_1_title(prs, fig_dir)
    slide_2_background(prs, fig_dir)
    slide_3_paradigm(prs, fig_dir)
    slide_4_results(prs, fig_dir)
    slide_5_model(prs, fig_dir)
    slide_6_takeaway(prs, fig_dir)

    prs.save(out_path)
    print(f"Saved: {out_path} ({os.path.getsize(out_path) / 1024:.0f} KB)")


if __name__ == "__main__":
    build(
        "blitz/output/zahno2026_test/tmp/figures",
        "blitz/output/zahno2026_test/paper_blitz_v3.pptx",
    )
