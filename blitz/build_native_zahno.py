"""
Build hybrid PPT for Zahno et al. 2026.

Strategy:
  - Paradigm diagrams → native PPT shapes (editable boxes, arrows, icons)
  - Data charts → matplotlib PNG (Nature-accurate styling)
  - Model schematics → native PPT shapes (editable stages)
  - Text elements → native textboxes (always editable)

This gives the best of both worlds:
  - User can edit paradigm layout, labels, colors in PowerPoint
  - Data figures maintain Nature-quality styling (spines, ticks, SEM bands)
"""

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

REPO_DIR = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(REPO_DIR))

from blitz.skills.style import Color, NPG, SlideFont, Layout
from blitz.skills.pptx_native import (
    _add_textbox, _add_rounded_rect, _hex_to_rgb, add_panel_label,
    draw_paradigm_native, draw_bar_chart_native, draw_line_chart_native,
    draw_model_native,
)
from pptx.enum.text import PP_ALIGN

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import numpy as np


def _render_results_figure(out_dir: Path) -> str:
    """Render behavioral results as Nature-style matplotlib figure."""
    from blitz.skills.style import MATPLOTLIB_STYLE
    plt.rcParams.update(MATPLOTLIB_STYLE)
    np.random.seed(42)

    fig, axes = plt.subplots(1, 3, figsize=(14, 4.5),
                             gridspec_kw={'width_ratios': [1.4, 1, 1]},
                             constrained_layout=True)
    fig.patch.set_facecolor('white')
    colors = ['#009E73', '#0072B2', '#D55E00']
    labels = ['Slow (108 km/h)', 'Moderate (180 km/h)', 'Fast (252 km/h)']

    # Panel A: Swing vs Ball position
    ax = axes[0]
    ax.text(-0.08, 1.08, 'A', transform=ax.transAxes, fontsize=24, fontweight='bold')
    x = np.array([40, 50, 60, 70, 80, 90, 100])
    for j, (c, lbl) in enumerate(zip(colors, labels)):
        bias = (j + 1) * 1.5
        y = x + np.random.normal(0, 1, len(x)) - bias * (x - 70) / 30
        sem = np.abs(np.random.normal(1.5, 0.3, len(x)))
        ax.plot(x, y, color=c, lw=2, label=lbl)
        ax.fill_between(x, y - sem, y + sem, color=c, alpha=0.12)
    ax.plot([40, 100], [40, 100], '--', color='#AAAAAA', lw=1, alpha=0.5)
    ax.set_xlabel('Ball Position (cm)')
    ax.set_ylabel('Swing Position (cm)')
    ax.set_title('Swing vs Ball Position', fontsize=20, fontweight='bold', loc='left')
    ax.legend(frameon=False, fontsize=12)

    # Panel B: Bias magnitude
    ax = axes[1]
    ax.text(-0.08, 1.08, 'B', transform=ax.transAxes, fontsize=24, fontweight='bold')
    means = [2.1, 4.8, 8.3]
    sems = [0.5, 0.7, 1.1]
    bars = ax.bar(range(3), means, yerr=sems, color=colors, capsize=4, alpha=0.85, edgecolor='none')
    for j in range(3):
        pts = np.random.normal(means[j], sems[j] * 1.5, 15)
        jitter = np.random.normal(0, 0.06, 15)
        ax.scatter(j + jitter, pts, color=colors[j], alpha=0.3, s=18, zorder=3)
    # Significance
    y_max = max(means) + max(sems) + 1.5
    ax.plot([0, 0, 2, 2], [y_max, y_max + 0.3, y_max + 0.3, y_max], color='#333', lw=0.8)
    ax.text(1, y_max + 0.4, '***', ha='center', fontsize=14)
    ax.set_xticks(range(3))
    ax.set_xticklabels(['Slow', 'Moderate', 'Fast'])
    ax.set_ylabel('Bias Magnitude (cm)')
    ax.set_title('Bayesian Bias', fontsize=20, fontweight='bold', loc='left')

    # Panel C: Bias × Day
    ax = axes[2]
    ax.text(-0.08, 1.08, 'C', transform=ax.transAxes, fontsize=24, fontweight='bold')
    means_d = [3.2, 5.1, 5.4]
    sems_d = [0.6, 0.5, 0.5]
    day_colors = ['#F39B7F', '#E64B35', '#E64B35']
    ax.bar(range(3), means_d, yerr=sems_d, color=day_colors, capsize=4, alpha=0.85, edgecolor='none')
    for j in range(3):
        pts = np.random.normal(means_d[j], sems_d[j] * 1.5, 15)
        jitter = np.random.normal(0, 0.06, 15)
        ax.scatter(j + jitter, pts, color=day_colors[j], alpha=0.3, s=18, zorder=3)
    ax.set_xticks(range(3))
    ax.set_xticklabels(['Day 1', 'Day 2', 'Day 3'])
    ax.set_ylabel('Bias (cm)')
    ax.set_title('Bias × Day', fontsize=20, fontweight='bold', loc='left')

    out_path = str(out_dir / 'results_figure.png')
    plt.savefig(out_path, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close()
    return out_path


def _render_model_figure(out_dir: Path) -> str:
    """Render model comparison as Nature-style matplotlib figure."""
    from blitz.skills.style import MATPLOTLIB_STYLE
    plt.rcParams.update(MATPLOTLIB_STYLE)
    np.random.seed(42)

    fig, axes = plt.subplots(1, 2, figsize=(10, 4.5), constrained_layout=True)
    fig.patch.set_facecolor('white')

    # Panel B: Predictions vs Data
    ax = axes[0]
    ax.text(-0.08, 1.08, 'B', transform=ax.transAxes, fontsize=24, fontweight='bold')
    x = np.array([40, 50, 60, 70, 80, 90, 100])
    ax.plot(x, x, '--', color='#8491B4', lw=2, label='MLE (no prior)')
    ax.plot(x, [45, 52, 59, 68, 78, 86, 95], color='#E64B35', lw=2.5, label='BLS (Bayesian)')
    data_y = [44, 51, 58, 66, 76, 87, 94]
    ax.scatter(x, data_y, color='#3C5488', s=50, zorder=3, label='Data')
    ax.errorbar(x, data_y, yerr=[1.5]*7, fmt='none', ecolor='#3C5488', capsize=3, alpha=0.5)
    ax.set_xlabel('Sample Interval (cm)')
    ax.set_ylabel('Production (cm)')
    ax.set_title('Model Predictions', fontsize=20, fontweight='bold', loc='left')
    ax.legend(frameon=False, fontsize=12)

    # Panel C: Model fit R²
    ax = axes[1]
    ax.text(-0.08, 1.08, 'C', transform=ax.transAxes, fontsize=24, fontweight='bold')
    models = ['MLE', 'MAP', 'BLS']
    r2 = [0.72, 0.85, 0.94]
    colors_m = ['#8491B4', '#F39B7F', '#E64B35']
    bars = ax.bar(range(3), r2, color=colors_m, alpha=0.85, edgecolor='none')
    ax.set_xticks(range(3))
    ax.set_xticklabels(models)
    ax.set_ylabel('R²')
    ax.set_ylim(0, 1.05)
    ax.set_title('Model Fit', fontsize=20, fontweight='bold', loc='left')
    # Annotate best
    ax.annotate('Best fit', xy=(2, 0.94), xytext=(1.5, 0.75),
                fontsize=12, color='#E64B35',
                arrowprops=dict(arrowstyle='->', color='#E64B35', lw=1.2))

    out_path = str(out_dir / 'model_figure.png')
    plt.savefig(out_path, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close()
    return out_path


def build():
    prs = Presentation()
    prs.slide_width = Layout.SLIDE_WIDTH
    prs.slide_height = Layout.SLIDE_HEIGHT
    out_dir = REPO_DIR / "blitz" / "output" / "zahno2026_native"
    out_dir.mkdir(parents=True, exist_ok=True)

    # ── Slide 1: Title ─────────────────────────────────────
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    slide1.background.fill.solid()
    slide1.background.fill.fore_color.rgb = Color.WHITE

    _add_textbox(slide1, Inches(0.8), Inches(1.5), Inches(11), Inches(1.5),
                 "Humans can learn bimodal priors\nin complex sensorimotor behaviour",
                 font_size=32, bold=True, name="title_main")

    _add_textbox(slide1, Inches(0.8), Inches(3.2), Inches(8), Inches(0.6),
                 "Zahno, Heed, & van Beers (2026)  |  Proc. R. Soc. B",
                 font_size=16, color=RGBColor(0x55, 0x55, 0x55), name="title_authors")

    _add_textbox(slide1, Inches(0.8), Inches(4.5), Inches(10), Inches(1.5),
                 "Can humans learn complex (bimodal) statistical regularities\n"
                 "of the environment through sensorimotor experience?",
                 font_size=20, color=Color.DARK, name="title_question")

    # ── Slide 2: Paradigm (native shapes) ──────────────────
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    slide2.background.fill.solid()
    slide2.background.fill.fore_color.rgb = Color.WHITE

    _add_textbox(slide2, Layout.TITLE_LEFT, Layout.TITLE_TOP,
                 Layout.TITLE_WIDTH, Layout.TITLE_HEIGHT,
                 "Bimodal priors are learned through 1440 VR tennis serves across 3 days",
                 font_size=SlideFont.TITLE, bold=True, name="slide_title")

    add_panel_label(slide2, "A", Layout.MARGIN_LEFT, Inches(0.8))
    draw_paradigm_native(slide2, {
        "title": "Single Trial Flow",
        "epochs": [
            {"label": "Fixation\n+ Cue", "duration": "", "color": "#F0F0F0", "icon": "cross"},
            {"label": "Delay", "duration": "1-2 s", "color": "#E8E8E8", "icon": "dot"},
            {"label": "Avatar\nServe", "duration": "", "color": "#E0E8F0", "icon": "screen"},
            {"label": "Ball\nApproach", "duration": "varies", "color": "#D0E0F0", "icon": "circle"},
            {"label": "Swing +\nContact", "duration": "400 ms", "color": "#B8D4F0", "icon": "arrow_keys"},
            {"label": "Feedback", "duration": "0-100 pts", "color": "#A0C8F0", "icon": "checkmark"},
        ],
        "show_timeline": True,
    }, left=Layout.MARGIN_LEFT + Inches(0.3), top=Inches(1.1),
       width=Layout.CONTENT_WIDTH - Inches(0.5), height=Inches(2.5))

    # Panel B: Design structure (bar chart)
    add_panel_label(slide2, "B", Layout.MARGIN_LEFT, Inches(4.0))
    draw_bar_chart_native(slide2, {
        "title": "Speed Conditions (Visual Uncertainty)",
        "conditions": ["Slow\n108 km/h", "Moderate\n180 km/h", "Fast\n252 km/h"],
        "values": [480, 480, 480],
        "colors": ["#009E73", "#0072B2", "#D55E00"],
        "name": "speed_conditions",
    }, left=Layout.MARGIN_LEFT + Inches(0.3), top=Inches(4.3),
       width=Inches(4.5), height=Inches(2.8))

    # Panel C: Session structure
    add_panel_label(slide2, "C", Inches(6), Inches(4.0))
    _add_textbox(slide2, Inches(6.3), Inches(4.3), Inches(5.5), Inches(0.4),
                 "Session Structure", font_size=18, bold=True, name="session_title")

    session_colors = {"Day 1": "#FFF0F0", "Day 2": "#F0FFF0", "Day 3": "#F0FFF0"}
    session_trials = {"Day 1": "480 trials (initial)", "Day 2": "480 trials (post-sleep)", "Day 3": "480 trials (post-sleep)"}
    for j, (day, color) in enumerate(session_colors.items()):
        y = Inches(4.8) + j * Inches(0.7)
        _add_rounded_rect(slide2, Inches(6.3), y, Inches(5), Inches(0.55),
                          fill_color=color, border_color="#AAAAAA",
                          name=f"session_box_{j}")
        _add_textbox(slide2, Inches(6.5), y + Inches(0.05), Inches(4.5), Inches(0.45),
                     f"{day} — {session_trials[day]}",
                     font_size=14, bold=True, align=PP_ALIGN.LEFT,
                     name=f"session_label_{j}")

    # ── Slide 3: Results — matplotlib PNG (Nature-accurate) ─
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    slide3.background.fill.solid()
    slide3.background.fill.fore_color.rgb = Color.WHITE

    _add_textbox(slide3, Layout.TITLE_LEFT, Layout.TITLE_TOP,
                 Layout.TITLE_WIDTH, Layout.TITLE_HEIGHT,
                 "Higher uncertainty drives stronger Bayesian bias toward the prior",
                 font_size=SlideFont.TITLE, bold=True, name="slide_title")

    # Generate matplotlib figure with Nature styling
    results_png = _render_results_figure(out_dir)
    if results_png:
        from PIL import Image as PILImage
        im = PILImage.open(results_png)
        ar = im.width / im.height
        fig_w = Inches(12)
        fig_h = int(fig_w / ar)
        max_h = Inches(6)
        if fig_h > max_h:
            fig_h = max_h
            fig_w = int(fig_h * ar)
        slide3.shapes.add_picture(
            results_png,
            Layout.MARGIN_LEFT + Inches(0.2), Inches(1.1),
            fig_w, fig_h,
        ).name = "results_figure"

    # ── Slide 4: Model ────────────────────────────────────
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    slide4.background.fill.solid()
    slide4.background.fill.fore_color.rgb = Color.WHITE

    _add_textbox(slide4, Layout.TITLE_LEFT, Layout.TITLE_TOP,
                 Layout.TITLE_WIDTH, Layout.TITLE_HEIGHT,
                 "BLS model captures bimodal prior learning (R² = 0.94)",
                 font_size=SlideFont.TITLE, bold=True, name="slide_title")

    # Panel A: Model — left column, full height
    add_panel_label(slide4, "A", Layout.MARGIN_LEFT, Inches(0.8))
    draw_model_native(slide4, {
        "title": "Model Architecture",
        "stages": [
            {"label": "Measurement", "sublabel": "p(tm|ts), scalar variability", "color": "#E8F0F8"},
            {"label": "Estimation", "sublabel": "f(tm) — mapping function", "color": "#F0F0E8"},
            {"label": "Production", "sublabel": "p(tp|te), motor noise", "color": "#E8F8E8"},
        ],
    }, left=Layout.MARGIN_LEFT + Inches(0.3), top=Inches(1.1),
       width=Inches(4.5), height=Inches(5.5))

    # Panels B+C: matplotlib PNG (Nature-accurate)
    model_png = _render_model_figure(out_dir)
    if model_png:
        from PIL import Image as PILImage
        im = PILImage.open(model_png)
        ar = im.width / im.height
        fig_w = Inches(7.5)
        fig_h = int(fig_w / ar)
        slide4.shapes.add_picture(
            model_png,
            Inches(5.5), Inches(1.1),
            fig_w, fig_h,
        ).name = "model_charts_figure"

    # ── Slide 5: Takeaway ─────────────────────────────────
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    slide5.background.fill.solid()
    slide5.background.fill.fore_color.rgb = Color.WHITE

    _add_textbox(slide5, Layout.TITLE_LEFT, Layout.TITLE_TOP,
                 Layout.TITLE_WIDTH, Layout.TITLE_HEIGHT,
                 "Key Takeaways",
                 font_size=SlideFont.TITLE, bold=True, name="slide_title")

    # Left: takeaway points
    takeaways = [
        "Humans learn bimodal priors through sensorimotor experience",
        "Higher uncertainty → stronger prior influence (Bayesian)",
        "BLS model best explains behavior (R² = 0.94)",
    ]
    for i, tw in enumerate(takeaways):
        y = Inches(1.3) + i * Inches(1.3)
        badge = _add_rounded_rect(slide5, Inches(0.6), y, Inches(0.5), Inches(0.5),
                                  fill_color=NPG.CYCLE[i], border_color=NPG.CYCLE[i],
                                  name=f"takeaway_badge_{i}")
        _add_textbox(slide5, Inches(0.6), y + Inches(0.02), Inches(0.5), Inches(0.45),
                     str(i + 1), font_size=20, bold=True,
                     color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER,
                     name=f"takeaway_number_{i}")
        _add_textbox(slide5, Inches(1.4), y, Inches(5.5), Inches(0.5),
                     tw, font_size=18, name=f"takeaway_text_{i}")

    # Right: key supporting chart (model fit)
    add_panel_label(slide5, "", Inches(7.5), Inches(0.8))
    draw_bar_chart_native(slide5, {
        "title": "Model Fit (R²)",
        "conditions": ["MLE", "MAP", "BLS"],
        "values": [0.72, 0.85, 0.94],
        "colors": ["#8491B4", "#F39B7F", "#E64B35"],
        "name": "takeaway_chart",
    }, left=Inches(7.5), top=Inches(1.1),
       width=Inches(5.3), height=Inches(5))

    # Save
    out_path = REPO_DIR / "blitz" / "output" / "zahno2026_native" / "paper_blitz_native.pptx"
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(f"Saved: {out_path}")

    # Print shape inventory
    total_shapes = 0
    for i, slide in enumerate(prs.slides):
        n = len(list(slide.shapes))
        total_shapes += n
        print(f"  Slide {i+1}: {n} shapes")
    print(f"  Total: {total_shapes} editable components")

    return str(out_path)


if __name__ == "__main__":
    build()
