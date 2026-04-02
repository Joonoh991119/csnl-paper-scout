"""
Hybrid PPT Builder — academic-pptx-skill philosophy + Paper Blitz patterns.

Philosophy (from academic-pptx-skill, strictly followed):
  - Action titles (complete sentence takeaways), not topic labels
  - Ghost deck test: titles alone tell the argument
  - One exhibit per slide, annotated with key finding
  - ~40 word body limit
  - Communication-first, not design-forward
  - Conclusions stay on screen for Q&A

Paper Blitz variations (slide patterns only):
  - Paper figure crops for ALL data figures (never AI-generated)
  - Native PPT shapes for paradigm diagrams (user-editable)
  - Native PPT shapes for model schematics (user-editable)
  - AI generation ONLY for: graphical abstract, paradigm, scenario prediction
  - Korean narration scripts (bilingual)

Design standards (from academic-pptx-skill):
  - Colors: primary #1F4E79 (navy), accent #2E75B6, body #2D2D2D, muted #777777
  - Font: Arial, single typeface throughout
  - Title: 24-28pt bold, Body: 20pt min, Chart labels: 16-18pt, Citations: 12-14pt
  - White background for content slides, dark navy for title/conclusions
"""

import json
import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

REPO_DIR = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(REPO_DIR))

from blitz.skills.style import Layout
from blitz.skills.pptx_native import (
    _add_textbox, _add_rounded_rect, add_panel_label,
    draw_paradigm_native, draw_model_native,
)
from blitz.skills.slide_spec import (
    ZONES, spec_paradigm_flow, spec_model_stages,
    validate_specs, spec_to_summary,
)
from blitz.skills.spec_renderer import render_specs

# ── Academic-PPTX-Skill Design Constants ───────────────────────

COLORS = {
    "bg": "FFFFFF",
    "primary": "1F4E79",    # dark navy — titles
    "accent": "2E75B6",     # mid-blue — headers, highlights
    "body": "2D2D2D",       # near-black — body text
    "muted": "777777",      # gray — citations, captions
    "rule": "CCCCCC",       # light gray — divider lines
    "highlight": "FFF2CC",  # yellow — callout boxes
}

FONTS = {
    "face": "Arial",
    "title": 26,
    "section_header": 22,
    "body": 20,
    "label": 16,
    "cite": 13,
}

MARGIN = Inches(0.5)

def _rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# ── Shared Elements ────────────────────────────────────────────

def _white_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(COLORS["bg"])


def _dark_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(COLORS["primary"])


def _action_title(slide, text: str):
    """Action title + thin divider rule. Academic-pptx-skill standard."""
    _add_textbox(slide, MARGIN, Inches(0.2), Inches(12), Inches(0.85),
                 text, font_size=FONTS["title"], bold=True,
                 color=_rgb(COLORS["primary"]), name="action_title",
                 font_name=FONTS["face"])
    # Thin divider under title
    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, MARGIN, Inches(1.05), Inches(12), Pt(2))
    rule.fill.solid()
    rule.fill.fore_color.rgb = _rgb(COLORS["rule"])
    rule.line.fill.background()
    rule.name = "title_rule"


def _citation(slide, text: str, y: float = 5.1):
    """In-slide citation at bottom."""
    _add_textbox(slide, MARGIN, Inches(y), Inches(12), Inches(0.3),
                 text, font_size=FONTS["cite"],
                 color=_rgb(COLORS["muted"]), name="citation",
                 font_name=FONTS["face"])


def _add_figure_crop(slide, path: str, left, top, max_w, max_h, name="figure"):
    """Add paper-cropped figure preserving aspect ratio."""
    if not path or not os.path.exists(path):
        return None
    from PIL import Image
    im = Image.open(path)
    ar = im.width / im.height
    w = max_w
    h = int(w / ar)
    if h > max_h:
        h = max_h
        w = int(h * ar)
    shape = slide.shapes.add_picture(path, left, top, w, h)
    shape.name = name
    return shape


def _resolve_figure(figure_id: str, figures: list, base_dir: Path) -> str | None:
    """Find figure file path by ID. Returns absolute path.

    Handles mismatches between LLM-assigned IDs (e.g., 'Fig. 3') and
    extracted IDs (e.g., 'p4_img0') by:
    1. Exact match on figure_id
    2. Fuzzy match: 'Fig. N' → find embedded_image on page N or near it
    3. Fallback: largest embedded_image
    """
    def _try_path(p: str) -> str | None:
        if not p:
            return None
        for candidate in [
            Path(p).resolve(),
            (base_dir / "figures" / os.path.basename(p)).resolve(),
            (Path.cwd() / p).resolve(),
        ]:
            if candidate.exists():
                return str(candidate)
        return None

    # 1. Exact match
    for fig in figures:
        if fig.get("figure_id") == figure_id:
            result = _try_path(fig.get("path", ""))
            if result:
                return result

    # 2. Fuzzy: extract number from "Fig. N" / "Figure N" and find best match
    import re
    m = re.search(r'(?:Fig\.?|Figure)\s*(\d+)', figure_id, re.IGNORECASE)
    if m:
        fig_num = int(m.group(1))

        # 2a. Best: figure_region with matching figure_number
        regions = [f for f in figures if f.get("type") == "figure_region"
                   and f.get("figure_number") == fig_num]
        if regions:
            result = _try_path(regions[0].get("path", ""))
            if result:
                return result

        # 2b. Any figure_region by index
        all_regions = sorted(
            [f for f in figures if f.get("type") == "figure_region"],
            key=lambda f: f.get("figure_number", 999))
        if fig_num - 1 < len(all_regions):
            result = _try_path(all_regions[fig_num - 1].get("path", ""))
            if result:
                return result

        # 2c. Embedded images
        embedded = [f for f in figures if f.get("type") == "embedded_image"]
        if embedded:
            embedded.sort(key=lambda f: abs(f.get("page", 0) - (fig_num + 1)))
            if fig_num - 1 < len(embedded):
                result = _try_path(embedded[fig_num - 1].get("path", ""))
                if result:
                    return result
            result = _try_path(embedded[0].get("path", ""))
            if result:
                return result

    # 3. Fallback: largest figure_region or embedded_image
    candidates = [f for f in figures if f.get("type") in ("figure_region", "embedded_image")]
    if candidates:
        candidates.sort(key=lambda f: f.get("area", 0), reverse=True)
        result = _try_path(candidates[0].get("path", ""))
        if result:
            return result

    # 4. page_full matching page number from figure_id (e.g., "p4_full")
    m2 = re.search(r'p(\d+)_full', figure_id)
    if m2:
        for fig in figures:
            if fig.get("figure_id") == figure_id:
                result = _try_path(fig.get("path", ""))
                if result:
                    return result

    # 5. Last resort: any full-page render (not page 1)
    for fig in figures:
        if fig.get("type") == "full_page" and fig.get("page", 0) > 1:
            result = _try_path(fig.get("path", ""))
            if result:
                return result

    return None


# ── Slide Builders ─────────────────────────────────────────────

def build_title(prs, data: dict, figures: list, base_dir: Path):
    """Title slide — dark navy bg, statement title. (academic-pptx-skill §1)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_bg(slide)

    # Main title
    _add_textbox(slide, Inches(0.7), Inches(1.4), Inches(11.5), Inches(1.8),
                 data.get("title_text", ""),
                 font_size=32, bold=True,
                 color=RGBColor(0xFF, 0xFF, 0xFF),
                 name="title_main", font_name=FONTS["face"])

    # Metadata (authors, journal, venue)
    meta_parts = []
    for el in data.get("elements", []):
        if el.get("style") == "metadata":
            meta_parts.append(el.get("content", ""))
    meta = "  ·  ".join(meta_parts) if meta_parts else ""
    if meta:
        # Thin accent rule
        rule = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(3.3), Inches(2.0), Pt(3))
        rule.fill.solid()
        rule.fill.fore_color.rgb = _rgb(COLORS["accent"])
        rule.line.fill.background()
        rule.name = "title_rule"

        _add_textbox(slide, Inches(0.7), Inches(3.5), Inches(11), Inches(0.5),
                     meta, font_size=15,
                     color=RGBColor(0xCA, 0xDC, 0xFC),
                     name="title_meta", font_name=FONTS["face"])

    # Why this paper (Paper Blitz specific)
    why = ""
    for el in data.get("elements", []):
        if el.get("style") == "subtitle" or "why" in el.get("style", "").lower():
            why = el.get("content", "")
            break
    if not why:
        why = data.get("why_this_paper", "")
    if why:
        _add_textbox(slide, Inches(0.7), Inches(4.2), Inches(10), Inches(1.0),
                     why, font_size=17,
                     color=RGBColor(0xDD, 0xDD, 0xDD),
                     name="title_why", font_name=FONTS["face"])


def build_motivation(prs, data: dict, figures: list, base_dir: Path):
    """Motivation/context slide. (academic-pptx-skill §2)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _white_bg(slide)
    _action_title(slide, data.get("title_text", ""))

    text_elements = [e for e in data.get("elements", []) if e.get("type") == "text_block"]
    y = Inches(1.2)
    for el in text_elements[:5]:
        content = el.get("content", "")
        _add_textbox(slide, MARGIN, y, Inches(12), Inches(0.5),
                     f"• {content}", font_size=FONTS["body"],
                     color=_rgb(COLORS["body"]),
                     name=f"motivation_bullet_{text_elements.index(el)}",
                     font_name=FONTS["face"])
        y += Inches(0.6)

    # Citation
    for el in data.get("elements", []):
        if el.get("style") == "citation":
            _citation(slide, el.get("content", ""))
            break


def build_methods(prs, data: dict, figures: list, base_dir: Path):
    """Methods slide — native paradigm diagram (Paper Blitz pattern).
    (academic-pptx-skill §4 + Paper Blitz native shapes)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _white_bg(slide)
    _action_title(slide, data.get("title_text", ""))

    # Check for native paradigm params
    paradigm = data.get("paradigm_params")
    fig_elements = [e for e in data.get("elements", []) if e.get("type") == "figure"]

    if paradigm:
        # Spec-based paradigm: compute exact coordinates, validate, render
        epochs = paradigm.get("epochs", [])
        specs = spec_paradigm_flow(epochs, ZONES["paradigm_flow"])
        issues = validate_specs(specs)
        if issues:
            for issue in issues:
                print(f"    [SPEC WARN] {issue}")
        render_specs(slide, specs)
    elif fig_elements:
        # Paper figure crop fallback
        fig_path = _resolve_figure(fig_elements[0].get("figure_id", ""), figures, base_dir)
        if fig_path:
            _add_figure_crop(slide, fig_path,
                             MARGIN, Inches(1.2), Inches(12), Inches(5.5),
                             name="methods_figure")

    # Text elements below paradigm (compact)
    text_elements = [e for e in data.get("elements", []) if e.get("type") == "text_block"]
    y = Inches(4.6)
    for el in text_elements[:2]:
        content = el.get("content", "")
        _add_textbox(slide, MARGIN, y, Inches(12), Inches(0.4),
                     content, font_size=FONTS["label"],
                     color=_rgb(COLORS["body"]),
                     name=f"methods_text_{text_elements.index(el)}",
                     font_name=FONTS["face"])
        y += Inches(0.45)


def build_results(prs, data: dict, figures: list, base_dir: Path):
    """Results slide — paper figure crop LEFT, interpretive text RIGHT.
    (academic-pptx-skill §5: figure left ~5.5", text right ~3.5")
    Also renders model_params if present (model specs on left, figure right)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _white_bg(slide)
    _action_title(slide, data.get("title_text", ""))

    # If this slide has model_params, render model schematic on left
    model_params = data.get("model_params")
    if model_params:
        stages = model_params.get("stages", [])
        if stages:
            specs = spec_model_stages(stages, ZONES["model_left"])
            issues = validate_specs(specs)
            if issues:
                for issue in issues:
                    print(f"    [SPEC WARN] {issue}")
            render_specs(slide, specs)

    fig_elements = [e for e in data.get("elements", []) if e.get("type") == "figure"]
    text_elements = [e for e in data.get("elements", []) if e.get("type") == "text_block"]

    if fig_elements:
        fig_path = _resolve_figure(fig_elements[0].get("figure_id", ""), figures, base_dir)

        if fig_path and text_elements:
            # Adjust layout when model schematic is on the left
            if model_params and model_params.get("stages"):
                fig_left = Inches(5.5)
                fig_max_w = Inches(7)
            else:
                fig_left = MARGIN
                fig_max_w = Inches(7.5)
            _add_figure_crop(slide, fig_path,
                             fig_left, Inches(1.2), fig_max_w, Inches(5.5),
                             name="results_figure")

            # Key finding annotation ON the figure (callout box)
            # Academic-pptx-skill: "annotate the key finding directly on the chart"
            key_finding = text_elements[0].get("content", "") if text_elements else ""

            if model_params and model_params.get("stages"):
                # Model schematic left + figure right: put key finding below figure
                if key_finding:
                    _add_textbox(slide, Inches(5.5), Inches(6.8), Inches(7), Inches(0.4),
                                 key_finding[:80], font_size=FONTS["cite"],
                                 color=_rgb(COLORS["muted"]),
                                 name="key_finding_text", font_name=FONTS["face"])
            else:
                # Standard layout: callout on figure + annotations right
                if key_finding:
                    callout = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(0.8), Inches(1.4), Inches(3.0), Inches(0.5))
                    callout.fill.solid()
                    callout.fill.fore_color.rgb = _rgb(COLORS["highlight"])
                    callout.line.color.rgb = _rgb("E6C800")
                    callout.line.width = Pt(1)
                    callout.name = "key_finding_callout"
                    _add_textbox(slide, Inches(0.85), Inches(1.42), Inches(2.9), Inches(0.46),
                                 key_finding[:60], font_size=14, bold=True,
                                 color=_rgb("7A5200"),
                                 align=PP_ALIGN.CENTER, name="key_finding_text",
                                 font_name=FONTS["face"])

                # Section header on right
                _add_textbox(slide, Inches(8.5), Inches(1.2), Inches(4), Inches(0.4),
                             "What to take away", font_size=FONTS["section_header"], bold=True,
                             color=_rgb(COLORS["accent"]),
                             name="results_header", font_name=FONTS["face"])

                y = Inches(1.7)
                for i, el in enumerate(text_elements[:3]):
                    _add_textbox(slide, Inches(8.5), y, Inches(4), Inches(0.8),
                                 f"• {el.get('content', '')}", font_size=FONTS["body"] - 1,
                                 color=_rgb(COLORS["body"]),
                                 name=f"results_annotation_{i}",
                                 font_name=FONTS["face"])
                    y += Inches(1.0)

        elif fig_path:
            # Full-width figure (no text annotations)
            _add_figure_crop(slide, fig_path,
                             MARGIN, Inches(1.2), Inches(12), Inches(5.8),
                             name="results_figure")

    # Citation
    cap = fig_elements[0].get("caption_label", "") if fig_elements else ""
    if cap:
        _citation(slide, cap)


def build_model(prs, data: dict, figures: list, base_dir: Path):
    """Model slide — native schematic LEFT + paper figure crop RIGHT.
    (Paper Blitz pattern: editable model architecture)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _white_bg(slide)
    _action_title(slide, data.get("title_text", ""))

    model_params = data.get("model_params")
    fig_elements = [e for e in data.get("elements", []) if e.get("type") == "figure"]

    if model_params:
        stages = model_params.get("stages", [])
        specs = spec_model_stages(stages, ZONES["model_left"])
        issues = validate_specs(specs)
        if issues:
            for issue in issues:
                print(f"    [SPEC WARN] {issue}")
        render_specs(slide, specs)

    # Model predictions/comparison figure (paper crop on right)
    if fig_elements:
        fig_path = _resolve_figure(fig_elements[0].get("figure_id", ""), figures, base_dir)
        if fig_path:
            left = Inches(5.5) if model_params else MARGIN
            width = Inches(7) if model_params else Inches(12)
            _add_figure_crop(slide, fig_path,
                             left, Inches(1.3), width, Inches(5.5),
                             name="model_figure")


def build_takeaway(prs, data: dict, figures: list, base_dir: Path):
    """Conclusions slide — dark navy, stays on screen for Q&A.
    (academic-pptx-skill §7: numbered takeaways, contact info)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_bg(slide)

    # "Conclusions" label
    _add_textbox(slide, MARGIN, Inches(0.25), Inches(12), Inches(0.4),
                 "Conclusions", font_size=20,
                 color=RGBColor(0xA0, 0xBB, 0xDD),
                 name="conclusions_label", font_name=FONTS["face"])

    # Accent rule
    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, MARGIN, Inches(0.7), Inches(12), Pt(3))
    rule.fill.solid()
    rule.fill.fore_color.rgb = _rgb(COLORS["accent"])
    rule.line.fill.background()

    # Numbered takeaways
    text_elements = [e for e in data.get("elements", []) if e.get("type") == "text_block"]
    y = Inches(0.9)
    for i, el in enumerate(text_elements[:4]):
        content = el.get("content", "")
        # Numbered format: "1. Bold key phrase: supporting detail"
        _add_textbox(slide, MARGIN, y, Inches(12), Inches(1.0),
                     f"{i+1}. {content}", font_size=FONTS["body"] + 1,
                     color=RGBColor(0xFF, 0xFF, 0xFF),
                     name=f"takeaway_{i}", font_name=FONTS["face"])
        y += Inches(1.1)

    # Contact (Paper Blitz: researcher info)
    _add_textbox(slide, MARGIN, Inches(5.0), Inches(10), Inches(0.4),
                 "Feedback welcome  |  CSNL lab meeting",
                 font_size=14, color=RGBColor(0xA0, 0xBB, 0xDD),
                 name="contact_info", font_name=FONTS["face"])


def build_references(prs, data: dict, figures: list, base_dir: Path):
    """References slide. (academic-pptx-skill §10: required)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _white_bg(slide)

    _add_textbox(slide, MARGIN, Inches(0.2), Inches(12), Inches(0.5),
                 "References", font_size=24, bold=True,
                 color=_rgb(COLORS["primary"]),
                 name="references_title", font_name=FONTS["face"])

    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, MARGIN, Inches(0.72), Inches(12), Pt(2))
    rule.fill.solid()
    rule.fill.fore_color.rgb = _rgb(COLORS["rule"])
    rule.line.fill.background()

    refs = data.get("references", [])
    if isinstance(refs, list):
        ref_text = "\n\n".join(refs[:15])
    else:
        ref_text = str(refs)

    _add_textbox(slide, MARGIN, Inches(0.85), Inches(12), Inches(6),
                 ref_text, font_size=FONTS["cite"],
                 color=_rgb(COLORS["body"]),
                 name="references_body", font_name=FONTS["face"])


# ── Main Builder ───────────────────────────────────────────────

SLIDE_BUILDERS = {
    "title": build_title,
    "background": build_results,     # background with figure uses results layout
    "motivation": build_motivation,  # text-only motivation
    "methods": build_methods,
    "results": build_results,
    "model": build_model,
    "takeaway": build_takeaway,
    "references": build_references,
}


def build_pptx_hybrid(plan: dict, figures: list, base_dir: Path, out_path: Path) -> Path:
    """Build hybrid PPTX following academic-pptx-skill philosophy.

    Paper crops for data figures + native shapes for paradigm/model.
    """
    prs = Presentation()
    prs.slide_width = Layout.SLIDE_WIDTH
    prs.slide_height = Layout.SLIDE_HEIGHT

    slides = plan.get("slides", [])
    for slide_data in slides:
        stype = slide_data.get("slide_type", "results")
        builder = SLIDE_BUILDERS.get(stype, build_results)
        builder(prs, slide_data, figures, base_dir)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))

    total = sum(len(list(s.shapes)) for s in prs.slides)
    print(f"  Hybrid PPTX: {len(prs.slides)} slides, {total} components")
    print(f"  Saved: {out_path}")

    return out_path
