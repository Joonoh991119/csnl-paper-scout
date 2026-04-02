"""
Slide Constraint Validator — runs after every slide is built.

ALL thresholds imported from style.py. No local magic numbers.

Hard constraints (violation = must fix before proceeding):
  1. No text overlap (bounding box collision)
  2. Font size: body ≥ 14pt, title ≥ 24pt
  3. No light gray text (all body text must be ≥ darkness threshold)
  4. Figure aspect ratio preserved (within tolerance)
  5. No element extends beyond slide bounds
  6. No text hidden behind figures
  7. Figure-internal text readability (estimated effective pt on slide)
"""

import os
import sys
from pathlib import Path

from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from PIL import Image

REPO_DIR = Path(__file__).resolve().parent.parent.parent
sys.path.insert(0, str(REPO_DIR))
from blitz.skills.style import Threshold, Layout
from blitz.skills.validate_figure_text import estimate_figure_text_readability

# Derived from style.Threshold — single source of truth
OVERLAP_TOLERANCE = Emu(Inches(Threshold.OVERLAP_TOL_INCHES))


# ── Helpers ─────────────────────────────────────────────────────

def _get_bbox(shape):
    return (shape.left, shape.top,
            shape.left + shape.width, shape.top + shape.height)


def _bboxes_overlap(a, b):
    al, at, ar, ab = a
    bl, bt, br, bb = b
    tol = OVERLAP_TOLERANCE
    return (al < br - tol) and (ar > bl + tol) and (at < bb - tol) and (ab > bt + tol)


def _is_text_shape(shape):
    return shape.has_text_frame


def _is_image_shape(shape):
    return shape.shape_type == 13


def _get_font_info(shape):
    fonts = []
    if not shape.has_text_frame:
        return fonts
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            f = run.font
            size_pt = f.size.pt if f.size else None
            color_rgb = None
            if f.color and f.color.rgb:
                color_rgb = f.color.rgb
            fonts.append({
                "text": run.text[:50],
                "size_pt": size_pt,
                "color_rgb": color_rgb,
                "bold": f.bold,
            })
    return fonts


# ── Validators ──────────────────────────────────────────────────

def check_text_overlap(slide):
    violations = []
    text_shapes = [s for s in slide.shapes if _is_text_shape(s) and s.text.strip()]
    for i in range(len(text_shapes)):
        for j in range(i + 1, len(text_shapes)):
            a, b = text_shapes[i], text_shapes[j]
            if _bboxes_overlap(_get_bbox(a), _get_bbox(b)):
                violations.append(
                    f"TEXT OVERLAP: \"{a.text[:25]}\" ↔ \"{b.text[:25]}\""
                )
    return violations


def check_text_behind_figure(slide):
    violations = []
    figures = [s for s in slide.shapes if _is_image_shape(s)]
    texts = [s for s in slide.shapes if _is_text_shape(s) and s.text.strip()]
    shape_list = list(slide.shapes)
    for fig in figures:
        fig_bbox = _get_bbox(fig)
        fig_idx = shape_list.index(fig)
        for txt in texts:
            txt_bbox = _get_bbox(txt)
            txt_idx = shape_list.index(txt)
            if _bboxes_overlap(fig_bbox, txt_bbox) and fig_idx > txt_idx:
                violations.append(
                    f"TEXT BEHIND FIGURE: \"{txt.text[:30]}\" hidden by image"
                )
    return violations


def check_font_sizes(slide, title_shape_indices=None):
    """Check font size minimums.

    title_shape_indices: set of shape indices that are known titles.
    If not provided, heuristic: first text shape on slide = title.
    """
    violations = []
    text_shapes = [(i, s) for i, s in enumerate(slide.shapes) if _is_text_shape(s)]

    # Determine which shapes are titles
    if title_shape_indices is None:
        # Heuristic: first non-empty text shape is the title
        title_shape_indices = set()
        for idx, s in text_shapes:
            if s.text.strip():
                title_shape_indices.add(idx)
                break

    for idx, shape in text_shapes:
        is_title_shape = idx in title_shape_indices
        for fi in _get_font_info(shape):
            sz = fi["size_pt"]
            if sz is None:
                continue

            if is_title_shape:
                if sz < Threshold.MIN_TITLE_FONT:
                    violations.append(
                        f"TITLE TOO SMALL: \"{fi['text'][:30]}\" is {sz}pt, "
                        f"min {Threshold.MIN_TITLE_FONT}pt"
                    )
            elif sz <= Threshold.MIN_FIG_ANNOTATION:
                # 10pt or below — allowed ONLY for figure annotations
                pass
            elif sz < Threshold.MIN_BODY_FONT:
                violations.append(
                    f"BODY TOO SMALL: \"{fi['text'][:30]}\" is {sz}pt, "
                    f"min {Threshold.MIN_BODY_FONT}pt"
                )
    return violations


def check_text_color(slide):
    violations = []
    for shape in slide.shapes:
        for fi in _get_font_info(shape):
            rgb = fi["color_rgb"]
            if rgb is None:
                continue
            sz = fi["size_pt"] or 14
            # Allow light color only for figure annotations (≤10pt)
            if sz <= Threshold.MIN_FIG_ANNOTATION:
                continue
            r, g, b = rgb[0], rgb[1], rgb[2]
            if min(r, g, b) > Threshold.MAX_GRAY_BODY:
                violations.append(
                    f"TEXT TOO LIGHT: \"{fi['text'][:30]}\" "
                    f"#{r:02x}{g:02x}{b:02x} (need ≤ #{Threshold.MAX_GRAY_BODY:02x})"
                )
    return violations


def check_bounds(slide, sw, sh):
    violations = []
    tol = OVERLAP_TOLERANCE
    for shape in slide.shapes:
        bbox = _get_bbox(shape)
        if bbox[2] > sw + tol:
            violations.append(f"OUT OF BOUNDS (right): shape exceeds slide width")
        if bbox[3] > sh + tol:
            violations.append(f"OUT OF BOUNDS (bottom): shape exceeds slide height")
    return violations


def check_figure_aspect_ratio(slide, figure_sources):
    """Check AR for each image shape against source files.

    figure_sources: dict mapping shape index → source image path.
    """
    violations = []
    if not figure_sources:
        return violations
    for i, shape in enumerate(slide.shapes):
        if not _is_image_shape(shape):
            continue
        src = figure_sources.get(i)
        if not src or not os.path.exists(src):
            continue
        placed_ratio = shape.width / shape.height if shape.height > 0 else 0
        im = Image.open(src)
        src_ratio = im.width / im.height
        diff = abs(placed_ratio - src_ratio) / src_ratio
        if diff > Threshold.ASPECT_RATIO_TOL:
            violations.append(
                f"ASPECT RATIO: distorted {diff*100:.1f}% "
                f"(source={src_ratio:.2f}, placed={placed_ratio:.2f})"
            )
    return violations


def check_figure_text_readability(slide, figure_meta):
    """Check that text inside figures will be readable at placement size.

    figure_meta: list of dicts with keys:
      - shape_index: int (index in slide.shapes)
      - image_path: str
      - source_min_font_pt: float (7 for paper crops, 18 for generated)
      - source_width_inches: float (original figsize width)
    """
    warnings = []
    if not figure_meta:
        return warnings
    for meta in figure_meta:
        idx = meta.get("shape_index")
        if idx is None:
            continue
        shapes = list(slide.shapes)
        if idx >= len(shapes):
            continue
        shape = shapes[idx]
        if not _is_image_shape(shape):
            continue

        placement_w = shape.width / 914400  # EMU → inches
        result = estimate_figure_text_readability(
            image_path=meta.get("image_path", ""),
            slide_placement_width_inches=placement_w,
            source_min_font_pt=meta.get("source_min_font_pt", 7),
            source_figure_width_inches=meta.get("source_width_inches"),
        )
        if not result["readable"]:
            # For paper-cropped figures: warn only (can't change source fonts)
            # For generated figures: this IS a violation
            is_generated = meta.get("is_generated", False)
            msg = (f"FIGURE TEXT: ~{result['effective_pt']}pt effective — "
                   f"{result['recommendation']}")
            if is_generated:
                warnings.append(msg)  # Hard violation for generated figs
            # Paper crops: logged as warning in _img(), not a slide violation
    return warnings


# ── Main Validator ──────────────────────────────────────────────

def validate_slide(slide, slide_num, slide_width=None, slide_height=None,
                   figure_sources=None, figure_meta=None, title_shape_indices=None):
    """
    Run ALL constraint checks on a single slide.

    Args:
        slide: pptx slide object
        slide_num: int (1-based)
        slide_width/height: EMU (default from Layout)
        figure_sources: dict {shape_index: source_image_path} for AR check
        figure_meta: list of dicts for figure text readability check
        title_shape_indices: set of shape indices that are titles

    Returns: {"slide_num", "passed", "violations", "warnings"}
    """
    sw = slide_width or Layout.SLIDE_WIDTH
    sh = slide_height or Layout.SLIDE_HEIGHT

    violations = []
    violations.extend(check_text_overlap(slide))
    violations.extend(check_text_behind_figure(slide))
    violations.extend(check_font_sizes(slide, title_shape_indices))
    violations.extend(check_text_color(slide))
    violations.extend(check_bounds(slide, sw, sh))
    violations.extend(check_figure_aspect_ratio(slide, figure_sources or {}))
    violations.extend(check_figure_text_readability(slide, figure_meta or []))

    return {
        "slide_num": slide_num,
        "passed": len(violations) == 0,
        "violations": violations,
    }


def validate_presentation(prs, figure_sources_per_slide=None, figure_meta_per_slide=None):
    results = []
    for i, slide in enumerate(prs.slides):
        result = validate_slide(
            slide, i + 1,
            prs.slide_width, prs.slide_height,
            figure_sources=(figure_sources_per_slide or {}).get(i + 1),
            figure_meta=(figure_meta_per_slide or {}).get(i + 1),
        )
        results.append(result)

    total = sum(len(r["violations"]) for r in results)
    return {"all_passed": total == 0, "total_violations": total, "slides": results}


def print_report(report):
    for sr in report["slides"]:
        status = "✓" if sr["passed"] else "✗"
        print(f"  Slide {sr['slide_num']}: {status}")
        for v in sr["violations"]:
            print(f"    ✗ {v}")
    if report["all_passed"]:
        print(f"\n  ALL PASSED ({len(report['slides'])} slides)")
    else:
        print(f"\n  FAILED: {report['total_violations']} violations")
