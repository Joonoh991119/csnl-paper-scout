"""
Spec Renderer — Render ElementSpecs to PPT shapes.

Takes validated coordinate specs and produces exact PPT shapes.
No spatial guessing — everything is pre-computed by slide_spec.py.
"""

import os
import sys
from pathlib import Path

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

REPO_DIR = Path(__file__).resolve().parent.parent.parent
sys.path.insert(0, str(REPO_DIR))

from blitz.skills.slide_spec import ElementSpec

FONT_FACE = "Arial"


def _rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def render_spec(slide, spec: ElementSpec):
    """Render one ElementSpec to a PPT shape."""
    p = spec.properties
    x, y, w, h = Inches(spec.x), Inches(spec.y), Inches(spec.w), Inches(spec.h)

    if spec.type == "rect":
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(p.get("fill", "#F0F0F0"))
        shape.line.color.rgb = _rgb(p.get("border", "#333333"))
        shape.line.width = Pt(p.get("border_width", 1.0))
        shape.text_frame.clear()
        shape.name = spec.id

    elif spec.type == "textbox":
        box = slide.shapes.add_textbox(x, y, w, h)
        box.name = spec.id
        tf = box.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        para.text = p.get("text", "")
        para.font.size = Pt(p.get("font_size", 16))
        para.font.bold = p.get("bold", False)
        para.font.color.rgb = _rgb(p.get("color", "#1A1A1A"))
        para.font.name = FONT_FACE
        align_map = {"center": PP_ALIGN.CENTER, "left": PP_ALIGN.LEFT, "right": PP_ALIGN.RIGHT}
        para.alignment = align_map.get(p.get("align", "center"), PP_ALIGN.CENTER)

    elif spec.type == "arrow":
        from pptx.enum.shapes import MSO_CONNECTOR_TYPE
        end_x = Inches(p.get("end_x", spec.x + spec.w))
        end_y = Inches(p.get("end_y", spec.y + spec.h))
        line = slide.shapes.add_connector(
            MSO_CONNECTOR_TYPE.STRAIGHT, x, y, end_x, end_y)
        line.line.color.rgb = _rgb(p.get("color", "#555555"))
        line.line.width = Pt(p.get("width", 1.2))
        line.name = spec.id

        # Arrowhead triangle
        arrow_size = Inches(0.12)
        if abs(spec.w) > abs(spec.h):
            tri = slide.shapes.add_shape(
                MSO_SHAPE.ISOSCELES_TRIANGLE,
                end_x - arrow_size, end_y - arrow_size // 2,
                arrow_size, arrow_size)
            tri.rotation = 90.0
        else:
            tri = slide.shapes.add_shape(
                MSO_SHAPE.ISOSCELES_TRIANGLE,
                end_x - arrow_size // 2, end_y - arrow_size,
                arrow_size, arrow_size)
            tri.rotation = 180.0
        tri.fill.solid()
        tri.fill.fore_color.rgb = _rgb(p.get("color", "#555555"))
        tri.line.fill.background()
        tri.name = f"{spec.id}_head"

    elif spec.type == "icon_strip":
        strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        strip.fill.solid()
        strip.fill.fore_color.rgb = _rgb(p.get("color", "#CC0000"))
        strip.line.fill.background()
        strip.name = spec.id

        symbol = p.get("symbol", "")
        if symbol:
            sym_box = slide.shapes.add_textbox(x, y - Inches(0.02), w, h + Inches(0.04))
            sym_box.name = f"{spec.id}_symbol"
            para = sym_box.text_frame.paragraphs[0]
            para.text = symbol
            para.font.size = Pt(14)
            para.font.bold = True
            para.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            para.font.name = FONT_FACE
            para.alignment = PP_ALIGN.CENTER

    elif spec.type == "figure":
        fig_path = p.get("path", "")
        if fig_path and os.path.exists(fig_path):
            from PIL import Image
            im = Image.open(fig_path)
            ar = im.width / im.height
            fw = w
            fh = int(fw / ar)
            if fh > h:
                fh = h
                fw = int(fh * ar)
            shape = slide.shapes.add_picture(fig_path, x, y, fw, fh)
            shape.name = spec.id


def render_specs(slide, specs: list[ElementSpec]):
    """Render all specs to a slide."""
    for spec in specs:
        render_spec(slide, spec)
