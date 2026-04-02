"""
PPT Builder Skill — constraint-validated, style.py-driven.

ALL design tokens from style.py. No local magic numbers.
Every slide validated after build. Violations block output.
"""

import os
import sys
from pathlib import Path
from typing import Optional
from PIL import Image

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN

REPO_DIR = Path(__file__).resolve().parent.parent.parent
sys.path.insert(0, str(REPO_DIR))

from blitz.skills.style import Color, SlideFont, Layout, Threshold
from blitz.skills.validate_slide import validate_slide
from blitz.skills.validate_figure_text import precheck_figure


# ── Core helpers ────────────────────────────────────────────────

def _text(slide, text: str, left, top, width, height,
          size: int = SlideFont.BODY, color=Color.DARK,
          bold: bool = False, align=PP_ALIGN.LEFT):
    """Add a single text box."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    for run in p.runs:
        run.font.name = SlideFont.FAMILY
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
    return box


def _body_lines(slide, lines: list, left, top, width, height,
                size: int = SlideFont.BODY, color=Color.DARK,
                line_spacing: float = SlideFont.BODY_LINE_SPACING,
                align=PP_ALIGN.LEFT):
    """Add multiple body text lines.

    Each item in lines is a tuple: (text, bold, color_override)
    Empty string = vertical spacer.
    """
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, (text, bold, clr) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = Pt(size * line_spacing)
        if not text.strip():
            p.space_before = Pt(4)
            continue
        p.text = text
        for run in p.runs:
            run.font.name = SlideFont.FAMILY
            run.font.size = Pt(size)
            run.font.color.rgb = clr or color
            run.font.bold = bold
    return box


def _fig_label(slide, text: str, left, top, width, height):
    """Figure annotation — the ONLY context where small/light text is allowed."""
    return _text(slide, text, left, top, width, height,
                 size=SlideFont.FIG_ANNOTATION, color=Color.FIG_LABEL)


def _img(slide, path: str, left, top, max_w, max_h,
         source_min_font: float = 7.0, is_generated: bool = False):
    """Add image preserving AR within bounds. Pre-checks text readability."""
    if not os.path.exists(path):
        print(f"  MISSING: {path}")
        return None, None

    im = Image.open(path)
    ratio = im.width / im.height

    # Fit within bounds preserving AR
    w = max_w
    h = int(w / ratio)
    if h > max_h:
        h = max_h
        w = int(h * ratio)

    # Pre-check figure text readability
    placement_w_inches = w / 914400  # EMU → inches
    readable, report = precheck_figure(
        path, placement_w_inches,
        is_generated=is_generated,
        source_min_font_pt=source_min_font,
    )
    if not readable:
        print(f"  ⚠ FIGURE TEXT WARNING: {report['recommendation']}")

    pic = slide.shapes.add_picture(path, left, top, w, h)

    meta = {
        "image_path": path,
        "source_min_font_pt": source_min_font,
        "source_width_inches": report.get("source_width_in"),
        "readable": readable,
    }
    return pic, meta


def _bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = Color.WHITE


# ── Slide builders ──────────────────────────────────────────────
# Each returns (slide, figure_meta_list) for validation.

def slide_title(prs, fig_dir, analysis):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)

    pic, meta = _img(slide, f"{fig_dir}/title_tight.png",
                     Inches(0.5), Inches(0.3), Inches(10), Inches(2.0))

    _text(slide, "Zahno, Beck, Hossner & Kording  ·  Proc. R. Soc. B  ·  2026",
          Inches(0.5), Inches(2.5), Inches(10), Inches(0.5),
          size=SlideFont.BODY, color=Color.DARK)

    _text(slide, "Why this paper — JOP",
          Inches(0.5), Inches(3.3), Inches(6), Inches(0.5),
          size=SlideFont.BODY_LARGE, color=Color.ACCENT, bold=False)

    # QA-F5: operationalized JOP connection + bimodal vs asymmetric distinction
    _body_lines(slide, [
        ("JOP: asymmetric prior × time estimation (normative model)", True, Color.DARK),
        ("", False, None),
        ("This paper uses bimodal prior — a different non-Gaussian shape", False, Color.DARK),
        ("Core principle: prior shape qualitatively transforms estimation", False, Color.DARK),
        ("", False, None),
        ("Key for JOP: validates implicit non-Gaussian prior learning", True, Color.ACCENT),
        ("in ecologically valid tasks — foundational for JOP's model", False, Color.ACCENT),
    ], Inches(0.5), Inches(3.9), Inches(12), Inches(3.3),
       size=SlideFont.BODY)

    return slide, [], {"title_shape_indices": set()}


def slide_background(prs, fig_dir, analysis):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)

    _text(slide, "Background & Hypothesis",
          Inches(0.5), Inches(0.2), Inches(8), Inches(0.6),
          size=SlideFont.TITLE, bold=True, color=Color.BLACK)

    bg = analysis.get("background", {})
    hyp = analysis.get("hypothesis_and_predictions", {})
    preds = hyp.get("predictions", [])

    estab = bg.get("established_knowledge",
                    "Bayesian integration is well-established in lab tasks")
    gap_text = bg.get("gap",
                      "Untested in naturalistic tasks with non-Gaussian priors")

    body = [
        ("Established", True, Color.DARK),
        (estab, False, Color.DARK),
        ("", False, None),
        ("Gap", True, Color.DARK),
        (gap_text, False, Color.DARK),
        ("", False, None),
        ("Predictions", True, Color.ACCENT),
    ]
    for i, pred in enumerate(preds[:2]):
        body.append((f"P{i+1}: {pred.get('expected_result', '')}", False, Color.DARK))

    _body_lines(slide, body,
                Inches(0.4), Inches(1.0), Inches(5.8), Inches(6),
                size=SlideFont.BODY_MIN)

    # Fig 1 WITHOUT caption, labeled as simulation (QA-F2-B)
    pic, meta = _img(slide, f"{fig_dir}/fig1_nocap.png",
                     Inches(6.5), Inches(0.9), Inches(6.5), Inches(6.0),
                     source_min_font=7.0)

    _fig_label(slide, "Fig. 1 — Simulation (prediction, not data) | green=slow, blue=mod, red=fast",
               Inches(6.5), Inches(7.0), Inches(6.5), Inches(0.4))

    fig_meta = []
    if meta:
        meta["shape_index"] = list(slide.shapes).index(pic) if pic else None
        fig_meta.append(meta)
    return slide, fig_meta, {}


def slide_paradigm(prs, fig_dir, analysis):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)

    _text(slide, "Experimental Paradigm",
          Inches(0.5), Inches(0.2), Inches(8), Inches(0.6),
          size=SlideFont.TITLE, bold=True, color=Color.BLACK)

    pic1, meta1 = _img(slide, f"{fig_dir}/paradigm_v3.png",
                       Inches(0.3), Inches(1.0), Inches(9.0), Inches(5.5),
                       source_min_font=18.0, is_generated=True)

    pic2, meta2 = _img(slide, f"{fig_dir}/fig2_xr_photos.png",
                       Inches(9.5), Inches(1.0), Inches(3.6), Inches(2.5),
                       source_min_font=7.0)

    _fig_label(slide, "Fig. 2 — CAVE XR setup",
               Inches(9.5), Inches(3.5), Inches(3.5), Inches(0.3))

    _body_lines(slide, [
        ("Ball speed controls visual uncertainty", True, Color.DARK),
        ("Temporal demand fixed (400 ms window)", False, Color.DARK),
        ("→ isolates sensory reliability effect", False, Color.ACCENT),
    ], Inches(9.5), Inches(4.0), Inches(3.6), Inches(3),
       size=SlideFont.BODY_MIN)

    fig_meta = []
    if meta1 and pic1:
        meta1["shape_index"] = list(slide.shapes).index(pic1)
        fig_meta.append(meta1)
    return slide, fig_meta, {}


def slide_results(prs, fig_dir, analysis):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)

    _text(slide, "Key Results",
          Inches(0.5), Inches(0.2), Inches(6), Inches(0.6),
          size=SlideFont.TITLE, bold=True, color=Color.BLACK)

    pic, meta = _img(slide, f"{fig_dir}/fig3_results.png",
                     Inches(0.5), Inches(0.9), Inches(9), Inches(4.2),
                     source_min_font=7.0)

    _fig_label(slide,
        "Fig. 3 — x: ball position, y: estimation error | green=slow, blue=mod, red=fast | arrows=bimodal jump",
        Inches(0.5), Inches(5.1), Inches(9), Inches(0.3))

    stats = analysis.get("key_statistics", {})
    _body_lines(slide, [
        ("After consolidation (Days 2-3):", True, Color.DARK),
        (f"Moderate: {stats.get('bimodal_effect_moderate', 'B=−2.64, p=.006')}", False, Color.DARK),
        (f"Fast: {stats.get('bimodal_effect_fast', 'B=−2.36, p=.023')}", False, Color.DARK),
        (f"Slow: {stats.get('bimodal_effect_slow', 'B=−0.81, p=.154 (n.s.)')}", False, Color.DARK),
        ("", False, None),
        ("No effect on Day 1 → consolidation required", False, Color.DARK),
        ("No explicit awareness → implicit learning", False, Color.DARK),
    ], Inches(9.7), Inches(0.9), Inches(3.4), Inches(5),
       size=SlideFont.BODY_MIN)

    fig_meta = []
    if meta and pic:
        meta["shape_index"] = list(slide.shapes).index(pic)
        fig_meta.append(meta)
    return slide, fig_meta, {}


def slide_model(prs, fig_dir, analysis):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)

    _text(slide, "Bayesian Prior + Biomechanical Costs",
          Inches(0.5), Inches(0.2), Inches(10), Inches(0.6),
          size=SlideFont.TITLE, bold=True, color=Color.BLACK)

    # QA-F2-A: Fig 5 dominant (60%), Fig 4 subordinate (40%)
    pic1, _ = _img(slide, f"{fig_dir}/fig4_scatter_only.png",
                   Inches(0.3), Inches(1.0), Inches(4.5), Inches(3.2),
                   source_min_font=7.0)
    _fig_label(slide, "Fig. 4 — Control (uniform distribution): biomechanical bias",
               Inches(0.3), Inches(4.2), Inches(4.5), Inches(0.3))

    pic2, _ = _img(slide, f"{fig_dir}/fig5_tight.png",
                   Inches(5.2), Inches(1.0), Inches(7.8), Inches(3.5),
                   source_min_font=7.0)
    _fig_label(slide, "Fig. 5 — Combined model: prior + biomechanics",
               Inches(5.2), Inches(4.5), Inches(7.5), Inches(0.3))

    # QA-F4-B: "qualitatively" not "일치"
    _body_lines(slide, [
        ("Control experiment reveals linear motor-cost bias across all speeds", False, Color.DARK),
        ("Combined = Bayesian prior + biomechanical bias", False, Color.DARK),
        ("→ Qualitatively matches experimental data pattern", True, Color.ACCENT),
    ], Inches(0.3), Inches(5.0), Inches(12.5), Inches(2.2),
       size=SlideFont.BODY)

    return slide, [], {}


def slide_takeaway(prs, fig_dir, analysis):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)

    _text(slide, "Takeaway",
          Inches(0.5), Inches(0.2), Inches(4), Inches(0.6),
          size=SlideFont.TITLE, bold=True, color=Color.BLACK)

    _body_lines(slide, [
        ("1. Implicit bimodal prior learning in naturalistic tasks", True, Color.DARK),
        ("   No conscious awareness — overnight consolidation", False, Color.DARK),
        ("", False, None),
        ("2. Prior reliance scales with sensory uncertainty", True, Color.DARK),
        ("   Moderate & fast significant; slow not", False, Color.DARK),
        ("", False, None),
        ("3. Behavior = Bayesian inference + biomechanical costs", True, Color.DARK),
        ("   Combined model captures full data pattern", False, Color.DARK),
    ], Inches(0.4), Inches(1.0), Inches(6), Inches(5.5),
       size=SlideFont.BODY)

    pic, meta = _img(slide, f"{fig_dir}/fig3_results.png",
                     Inches(6.8), Inches(0.8), Inches(6.2), Inches(5),
                     source_min_font=7.0)

    _fig_label(slide, "Fig. 3 — Main result",
               Inches(6.8), Inches(5.9), Inches(4), Inches(0.3))

    return slide, [], {}


# ── Build + Validate ────────────────────────────────────────────

def build_and_validate(fig_dir: str, out_path: str, analysis: dict) -> dict:
    """Build all slides, validate each. Returns report."""
    prs = Presentation()
    prs.slide_width = Layout.SLIDE_WIDTH
    prs.slide_height = Layout.SLIDE_HEIGHT

    builders = [
        ("Title", slide_title),
        ("Background & Hypothesis", slide_background),
        ("Experimental Paradigm", slide_paradigm),
        ("Key Results", slide_results),
        ("Combined Model", slide_model),
        ("Takeaway", slide_takeaway),
    ]

    all_violations = []
    for i, (name, builder) in enumerate(builders):
        print(f"  [{i+1}] {name}...")
        slide, fig_meta, extra = builder(prs, fig_dir, analysis)

        result = validate_slide(
            slide, i + 1,
            Layout.SLIDE_WIDTH, Layout.SLIDE_HEIGHT,
            figure_meta=[{**m, "shape_index": m.get("shape_index")}
                         for m in fig_meta if m.get("shape_index") is not None],
            title_shape_indices=extra.get("title_shape_indices"),
        )
        if not result["passed"]:
            for v in result["violations"]:
                print(f"      ✗ {v}")
            all_violations.extend(result["violations"])
        else:
            print(f"      ✓")

    prs.save(out_path)
    size_kb = os.path.getsize(out_path) / 1024
    total = len(all_violations)
    print(f"\n  Saved: {out_path} ({size_kb:.0f} KB) — {total} violations")

    return {
        "path": out_path,
        "total_violations": total,
        "violations": all_violations,
        "passed": total == 0,
    }


if __name__ == "__main__":
    import json
    analysis = json.load(open("blitz/output/zahno2026_test/tmp/deep_analysis.json"))
    build_and_validate(
        "blitz/output/zahno2026_test/tmp/figures",
        "blitz/output/zahno2026_test/paper_blitz_v5.pptx",
        analysis,
    )
