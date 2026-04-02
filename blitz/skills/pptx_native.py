"""
PPT Native Renderer — Figures as editable PowerPoint shapes.

Instead of rendering matplotlib → PNG → insert image,
this module draws figures directly using python-pptx shapes:
  - Rectangles, rounded rectangles for boxes
  - Connectors for arrows
  - TextBoxes for labels
  - Native charts for data plots
  - Grouped shapes for semantic components

Every element gets a tagged name (e.g., "paradigm_epoch_2", "bar_chart_cond_A")
so the feedback collector can track user modifications.
"""

import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

REPO_DIR = Path(__file__).resolve().parent.parent.parent
sys.path.insert(0, str(REPO_DIR))

from blitz.skills.style import (
    Color, NPG, SlideFont, Layout, ParadigmColor, ParadigmLayout,
)

# ── Helpers ────────────────────────────────────────────────────

def _hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _add_textbox(slide, left, top, width, height, text, font_size=16,
                 bold=False, color=Color.BLACK, align=PP_ALIGN.LEFT,
                 name=None, font_name="Helvetica Neue"):
    """Add a tagged textbox."""
    box = slide.shapes.add_textbox(left, top, width, height)
    if name:
        box.name = name
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return box


def _add_rounded_rect(slide, left, top, width, height,
                      fill_color="#F0F0F0", border_color="#333333",
                      border_width=Pt(1), name=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgb(fill_color)
    shape.line.color.rgb = _hex_to_rgb(border_color)
    shape.line.width = border_width
    if name:
        shape.name = name
    # Remove default text
    shape.text_frame.clear()
    return shape


def _add_arrow(slide, start_x, start_y, end_x, end_y,
               color="#555555", width=Pt(1.2), name=None):
    """Add an arrow using line + triangle arrowhead shape."""
    # Line body
    line = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        start_x, start_y, end_x, end_y,
    )
    line.line.color.rgb = _hex_to_rgb(color)
    line.line.width = width
    if name:
        line.name = name

    # Arrowhead: small right-pointing triangle at end
    arrow_size = Inches(0.12)
    if abs(end_x - start_x) > abs(end_y - start_y):
        # Horizontal arrow → right-pointing triangle
        tri = slide.shapes.add_shape(
            MSO_SHAPE.ISOSCELES_TRIANGLE,
            end_x - arrow_size, end_y - arrow_size // 2,
            arrow_size, arrow_size,
        )
        tri.rotation = 90.0  # Point right
    else:
        # Vertical arrow → down-pointing triangle
        tri = slide.shapes.add_shape(
            MSO_SHAPE.ISOSCELES_TRIANGLE,
            end_x - arrow_size // 2, end_y - arrow_size,
            arrow_size, arrow_size,
        )
        tri.rotation = 180.0  # Point down

    tri.fill.solid()
    tri.fill.fore_color.rgb = _hex_to_rgb(color)
    tri.line.fill.background()
    tri.name = f"{name}_head" if name else "arrow_head"

    return line


def _add_line(slide, start_x, start_y, end_x, end_y,
              color="#AAAAAA", width=Pt(1.5), name=None):
    """Add a plain line (no arrowhead)."""
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        start_x, start_y, end_x, end_y,
    )
    connector.line.color.rgb = _hex_to_rgb(color)
    connector.line.width = width
    if name:
        connector.name = name
    return connector


# ── Paradigm Diagram (Native PPT) ─────────────────────────────

def draw_paradigm_native(slide, params: dict,
                         left: int = None, top: int = None,
                         width: int = None, height: int = None):
    """Draw paradigm diagram as native PPT shapes.

    Each epoch = rounded rect + label textbox + duration textbox + icon
    Arrows between epochs = connectors
    Timeline = line at bottom

    All elements tagged with names for feedback tracking.

    params schema (same as style_templates.py):
    {
        "title": str,
        "epochs": [{"label", "duration", "color", "icon"}, ...],
        "show_timeline": bool,
    }
    """
    left = left or Layout.MARGIN_LEFT
    top = top or Layout.CONTENT_TOP
    width = width or Layout.CONTENT_WIDTH
    height = height or Inches(3.0)

    epochs = params.get("epochs", [])
    title = params.get("title", "")
    n = len(epochs)
    if n == 0:
        return

    # Title
    if title:
        _add_textbox(slide, left, top - Inches(0.5), width, Inches(0.5),
                     title, font_size=22, bold=True, name="paradigm_title")

    # Calculate box dimensions
    box_w = int(width * 0.75 / n)  # Leave gaps
    box_h = Inches(1.2)
    gap = int(width * 0.25 / max(n - 1, 1)) if n > 1 else 0
    y = top + Inches(0.3)

    for i, ep in enumerate(epochs):
        x = left + i * (box_w + gap)
        fill = ep.get("color", "#F0F0F0")
        label = ep.get("label", f"Epoch {i+1}")
        duration = ep.get("duration", "")
        icon = ep.get("icon", "none")

        # Box
        box = _add_rounded_rect(
            slide, x, y, box_w, box_h,
            fill_color=fill, border_color="#333333",
            name=f"paradigm_epoch_{i}",
        )

        # Icon indicator (small colored strip at top of box)
        icon_colors = {
            "cross": "#CC0000", "grating": "#3C5488", "dot": "#333333",
            "arrow_keys": "#00A087", "checkmark": "#009E73",
            "question": "#888888", "screen": "#4DBBD5",
        }
        if icon in icon_colors:
            strip_h = Inches(0.15)
            strip = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x, y, box_w, strip_h,
            )
            strip.fill.solid()
            strip.fill.fore_color.rgb = _hex_to_rgb(icon_colors[icon])
            strip.line.fill.background()
            strip.name = f"paradigm_icon_{i}"

            # Icon text symbol inside strip
            icon_symbols = {
                "cross": "+", "grating": "≡", "dot": "●",
                "arrow_keys": "↕", "checkmark": "✓",
                "question": "?", "screen": "▢",
            }
            if icon in icon_symbols:
                _add_textbox(
                    slide, x, y - Inches(0.02), box_w, strip_h + Inches(0.04),
                    icon_symbols[icon], font_size=14, bold=True,
                    color=RGBColor(0xFF, 0xFF, 0xFF),
                    align=PP_ALIGN.CENTER,
                    name=f"paradigm_icon_symbol_{i}",
                )

        # Label (center of box)
        _add_textbox(
            slide, x + Inches(0.05), y + Inches(0.25),
            box_w - Inches(0.1), Inches(0.5),
            label, font_size=16, bold=True,
            align=PP_ALIGN.CENTER,
            name=f"paradigm_label_{i}",
        )

        # Duration (below label)
        if duration:
            _add_textbox(
                slide, x + Inches(0.05), y + box_h - Inches(0.35),
                box_w - Inches(0.1), Inches(0.3),
                duration, font_size=12, bold=False,
                color=RGBColor(0x55, 0x55, 0x55),
                align=PP_ALIGN.CENTER,
                name=f"paradigm_duration_{i}",
            )

        # Arrow to next epoch
        if i < n - 1:
            arrow_start_x = x + box_w
            arrow_end_x = x + box_w + gap
            arrow_y = y + box_h // 2
            _add_arrow(
                slide, arrow_start_x, arrow_y, arrow_end_x, arrow_y,
                color="#555555", name=f"paradigm_arrow_{i}",
            )

    # Timeline
    if params.get("show_timeline", True):
        tl_y = y + box_h + Inches(0.3)
        tl_start = left
        tl_end = left + (n - 1) * (box_w + gap) + box_w
        _add_line(slide, tl_start, tl_y, tl_end, tl_y,
                  color="#AAAAAA", width=Pt(1.5), name="paradigm_timeline")
        _add_textbox(
            slide, int((tl_start + tl_end) / 2) - Inches(0.3), tl_y + Inches(0.05),
            Inches(0.6), Inches(0.3),
            "Time →", font_size=11, color=RGBColor(0x88, 0x88, 0x88),
            align=PP_ALIGN.CENTER, name="paradigm_timeline_label",
        )


# ── Bar Chart (Native PPT Chart) ──────────────────────────────

def draw_bar_chart_native(slide, params: dict,
                          left: int = None, top: int = None,
                          width: int = None, height: int = None):
    """Draw a bar chart as native PPT chart object.

    Fully editable in PowerPoint — user can change data, colors, labels.

    params:
    {
        "title": str,
        "conditions": ["Low", "Med", "High"],
        "values": [4.2, 6.1, 7.8],
        "errors": [0.5, 0.4, 0.6],  # optional
        "colors": ["#E64B35", "#4DBBD5", "#00A087"],  # optional
        "x_label": str,
        "y_label": str,
    }
    """
    left = left or Inches(1)
    top = top or Inches(2)
    width = width or Inches(5)
    height = height or Inches(3.5)

    conditions = params.get("conditions", ["A", "B", "C"])
    values = params.get("values", [5, 6, 7])
    colors = params.get("colors", NPG.CYCLE[:len(conditions)])

    chart_data = CategoryChartData()
    chart_data.categories = conditions
    chart_data.add_series("Data", values)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        left, top, width, height,
        chart_data,
    )
    chart_frame.name = params.get("name", "bar_chart")

    chart = chart_frame.chart
    chart.has_legend = False

    # Style: remove top/right border, set colors
    plot = chart.plots[0]
    plot.gap_width = 80

    series = plot.series[0]
    for i, point in enumerate(series.points):
        point.format.fill.solid()
        c = colors[i % len(colors)] if i < len(colors) else colors[0]
        point.format.fill.fore_color.rgb = _hex_to_rgb(c)

    # Axis styling
    cat_axis = chart.category_axis
    cat_axis.has_minor_gridlines = False
    cat_axis.has_major_gridlines = False
    cat_axis.tick_labels.font.size = Pt(14)
    cat_axis.tick_labels.font.name = "Helvetica Neue"

    val_axis = chart.value_axis
    val_axis.has_minor_gridlines = False
    val_axis.has_major_gridlines = False
    val_axis.tick_labels.font.size = Pt(14)
    val_axis.tick_labels.font.name = "Helvetica Neue"

    # Title
    chart_title = params.get("title", "")
    if chart_title:
        chart.has_title = True
        chart.chart_title.text_frame.text = chart_title
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(18)
        chart.chart_title.text_frame.paragraphs[0].font.bold = True
        chart.chart_title.text_frame.paragraphs[0].font.name = "Helvetica Neue"

    return chart_frame


# ── Line Chart (Native PPT Chart) ─────────────────────────────

def draw_line_chart_native(slide, params: dict,
                           left: int = None, top: int = None,
                           width: int = None, height: int = None):
    """Draw a line chart as native PPT chart object.

    params:
    {
        "title": str,
        "categories": ["0", "1", "2", ...],  # x-axis values
        "series": [
            {"name": "Cond 1", "values": [1.2, 3.4, ...], "color": "#E64B35"},
            ...
        ],
    }
    """
    left = left or Inches(1)
    top = top or Inches(2)
    width = width or Inches(5)
    height = height or Inches(3.5)

    series_list = params.get("series", [])
    categories = params.get("categories", [str(i) for i in range(10)])

    chart_data = CategoryChartData()
    chart_data.categories = categories
    for s in series_list:
        chart_data.add_series(s.get("name", ""), s.get("values", []))

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE,
        left, top, width, height,
        chart_data,
    )
    chart_frame.name = params.get("name", "line_chart")

    chart = chart_frame.chart

    # Style series colors
    plot = chart.plots[0]
    for i, series in enumerate(plot.series):
        color = series_list[i].get("color", NPG.CYCLE[i % len(NPG.CYCLE)]) if i < len(series_list) else NPG.CYCLE[i % len(NPG.CYCLE)]
        series.format.line.color.rgb = _hex_to_rgb(color)
        series.format.line.width = Pt(2.0)
        series.smooth = True

    # Legend
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(12)

    # Axes
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(12)
    val_axis = chart.value_axis
    val_axis.has_major_gridlines = False
    val_axis.tick_labels.font.size = Pt(12)

    # Title
    chart_title = params.get("title", "")
    if chart_title:
        chart.has_title = True
        chart.chart_title.text_frame.text = chart_title
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(18)
        chart.chart_title.text_frame.paragraphs[0].font.bold = True

    return chart_frame


# ── Model/Schematic Diagram (Native PPT) ──────────────────────

def draw_model_native(slide, params: dict,
                      left: int = None, top: int = None,
                      width: int = None, height: int = None):
    """Draw model schematic as native PPT shapes.

    Vertical stage layout with arrows between stages.

    params:
    {
        "title": str,
        "stages": [{"label", "sublabel", "color"}, ...],
        "feedback": bool,
    }
    """
    left = left or Inches(0.5)
    top = top or Layout.CONTENT_TOP
    width = width or Inches(5)
    height = height or Inches(5)

    stages = params.get("stages", [])
    title = params.get("title", "")
    n = len(stages)
    if n == 0:
        return

    if title:
        _add_textbox(slide, left, top - Inches(0.4), width, Inches(0.4),
                     title, font_size=20, bold=True, name="model_title")

    box_w = Inches(3)
    box_h = Inches(0.7)
    gap_y = Inches(0.4)
    x_center = left + (width - box_w) // 2

    for i, stage in enumerate(stages):
        y = top + i * (box_h + gap_y)
        fill = stage.get("color", "#F0F0F0")

        # Box
        _add_rounded_rect(
            slide, x_center, y, box_w, box_h,
            fill_color=fill, name=f"model_stage_{i}",
        )

        # Label
        _add_textbox(
            slide, x_center + Inches(0.1), y + Inches(0.08),
            box_w - Inches(0.2), Inches(0.35),
            stage.get("label", ""), font_size=16, bold=True,
            align=PP_ALIGN.CENTER, name=f"model_label_{i}",
        )

        # Sublabel
        sublabel = stage.get("sublabel", "")
        if sublabel:
            _add_textbox(
                slide, x_center + Inches(0.1), y + Inches(0.38),
                box_w - Inches(0.2), Inches(0.25),
                sublabel, font_size=12, bold=False,
                color=RGBColor(0x55, 0x55, 0x55),
                align=PP_ALIGN.CENTER, name=f"model_sublabel_{i}",
            )

        # Arrow to next stage
        if i < n - 1:
            arrow_x = x_center + box_w // 2
            _add_arrow(
                slide, arrow_x, y + box_h,
                arrow_x, y + box_h + gap_y,
                name=f"model_arrow_{i}",
            )


# ── Panel Label ────────────────────────────────────────────────

def add_panel_label(slide, label: str, left: int, top: int, name: str = None):
    """Add a bold uppercase panel label (A, B, C...)."""
    _add_textbox(
        slide, left, top, Inches(0.4), Inches(0.4),
        label, font_size=24, bold=True,
        name=name or f"panel_label_{label}",
    )


# ── Composite: Full Figure Slide ───────────────────────────────

def build_paradigm_slide(prs, params: dict) -> None:
    """Build a complete paradigm slide with native shapes.

    params:
    {
        "slide_title": str,
        "paradigm": {epochs, title, ...},
        "bottom_chart": {type: "bar"|"line", ...} | null,
    }
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Background
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = Color.WHITE

    # Slide title
    _add_textbox(
        slide, Layout.TITLE_LEFT, Layout.TITLE_TOP,
        Layout.TITLE_WIDTH, Layout.TITLE_HEIGHT,
        params.get("slide_title", ""),
        font_size=SlideFont.TITLE, bold=True,
        name="slide_title",
    )

    # Panel A: Paradigm
    add_panel_label(slide, "A", Layout.MARGIN_LEFT, Layout.CONTENT_TOP - Inches(0.3))
    draw_paradigm_native(
        slide, params.get("paradigm", {}),
        left=Layout.MARGIN_LEFT + Inches(0.4),
        top=Layout.CONTENT_TOP,
        width=Layout.CONTENT_WIDTH - Inches(0.5),
        height=Inches(2.5),
    )

    # Panel B: Chart (if provided)
    bottom_chart = params.get("bottom_chart")
    if bottom_chart:
        chart_top = Layout.CONTENT_TOP + Inches(3.2)
        add_panel_label(slide, "B", Layout.MARGIN_LEFT, chart_top - Inches(0.3))

        chart_type = bottom_chart.get("type", "bar")
        if chart_type == "bar":
            draw_bar_chart_native(
                slide, bottom_chart,
                left=Layout.MARGIN_LEFT + Inches(0.5),
                top=chart_top,
                width=Inches(5), height=Inches(3),
            )
        elif chart_type == "line":
            draw_line_chart_native(
                slide, bottom_chart,
                left=Layout.MARGIN_LEFT + Inches(0.5),
                top=chart_top,
                width=Inches(5), height=Inches(3),
            )


def build_results_slide(prs, params: dict) -> None:
    """Build a results slide with native chart + annotations."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = Color.WHITE

    _add_textbox(
        slide, Layout.TITLE_LEFT, Layout.TITLE_TOP,
        Layout.TITLE_WIDTH, Layout.TITLE_HEIGHT,
        params.get("slide_title", "Results"),
        font_size=SlideFont.TITLE, bold=True,
        name="slide_title",
    )

    charts = params.get("charts", [])
    n_charts = len(charts)
    chart_w = int((Layout.CONTENT_WIDTH - Inches(0.5) * max(n_charts - 1, 0)) / max(n_charts, 1))

    for i, chart_params in enumerate(charts):
        x = Layout.MARGIN_LEFT + i * (chart_w + Inches(0.5))
        y = Layout.CONTENT_TOP

        add_panel_label(slide, chr(65 + i), x, y - Inches(0.3))

        ct = chart_params.get("type", "bar")
        if ct == "bar":
            draw_bar_chart_native(slide, chart_params, x + Inches(0.3), y, chart_w, Inches(4.5))
        elif ct == "line":
            draw_line_chart_native(slide, chart_params, x + Inches(0.3), y, chart_w, Inches(4.5))


# ── CLI Test ───────────────────────────────────────────────────

if __name__ == "__main__":
    prs = Presentation()
    prs.slide_width = Layout.SLIDE_WIDTH
    prs.slide_height = Layout.SLIDE_HEIGHT

    # Test paradigm slide
    build_paradigm_slide(prs, {
        "slide_title": "Experimental Design",
        "paradigm": {
            "title": "Single Trial Flow",
            "epochs": [
                {"label": "Fixation", "duration": "500ms", "color": "#F0F0F0", "icon": "cross"},
                {"label": "Stimulus", "duration": "200ms", "color": "#E0E8F0", "icon": "grating"},
                {"label": "Delay", "duration": "1-2s", "color": "#E8E8E8", "icon": "dot"},
                {"label": "Response", "duration": "until", "color": "#D8E8D8", "icon": "arrow_keys"},
                {"label": "Feedback", "duration": "500ms", "color": "#E8D8E8", "icon": "checkmark"},
            ],
            "show_timeline": True,
        },
        "bottom_chart": {
            "type": "bar",
            "title": "Accuracy by Condition",
            "conditions": ["Easy", "Medium", "Hard"],
            "values": [92, 78, 61],
            "colors": ["#00A087", "#4DBBD5", "#E64B35"],
            "name": "accuracy_chart",
        },
    })

    # Test results slide
    build_results_slide(prs, {
        "slide_title": "Behavioral Results",
        "charts": [
            {
                "type": "bar",
                "title": "Response Time",
                "conditions": ["Short", "Medium", "Long"],
                "values": [450, 520, 680],
                "colors": ["#E64B35", "#4DBBD5", "#00A087"],
                "name": "rt_chart",
            },
            {
                "type": "line",
                "title": "Learning Curve",
                "categories": [str(i) for i in range(1, 11)],
                "series": [
                    {"name": "Easy", "values": [60, 65, 72, 78, 82, 85, 87, 89, 90, 91], "color": "#00A087"},
                    {"name": "Hard", "values": [50, 52, 55, 58, 62, 65, 68, 70, 72, 74], "color": "#E64B35"},
                ],
                "name": "learning_chart",
            },
        ],
    })

    out = REPO_DIR / "blitz" / "style_knowledge" / "template_test" / "native_test.pptx"
    prs.save(str(out))
    print(f"Saved: {out}")

    # Print all shape names for verification
    for i, slide in enumerate(prs.slides):
        print(f"\nSlide {i+1}:")
        for shape in slide.shapes:
            print(f"  {shape.name:40s} type={shape.shape_type}")
