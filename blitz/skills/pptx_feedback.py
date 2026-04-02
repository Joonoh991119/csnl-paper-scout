"""
PPT Feedback Collector — Extract user modifications and update style tokens.

Workflow:
  1. Original PPT generated with tagged shapes (pptx_native.py)
  2. User opens in PowerPoint, modifies (colors, sizes, positions, text)
  3. This module reads both original + modified, computes diff
  4. Diffs are translated to style token updates

Tracked modifications:
  - Color changes (fill, line) → palette preferences
  - Size changes (width, height) → layout proportions
  - Position changes → spacing preferences
  - Text changes → label preferences
  - Deleted shapes → unwanted elements
  - Added shapes → missing elements
"""

import json
import sys
from datetime import datetime, timezone
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu

REPO_DIR = Path(__file__).resolve().parent.parent.parent
sys.path.insert(0, str(REPO_DIR))

FEEDBACK_DIR = REPO_DIR / "blitz" / "style_knowledge" / "feedback"


# ── Shape Extraction ──────────────────────────────────────────

def _extract_shape_data(shape) -> dict:
    """Extract all editable properties from a shape."""
    data = {
        "name": shape.name,
        "shape_type": str(shape.shape_type),
        "left": shape.left,
        "top": shape.top,
        "width": shape.width,
        "height": shape.height,
    }

    # Fill color
    try:
        if shape.fill and shape.fill.type is not None:
            fc = shape.fill.fore_color
            if fc and fc.rgb:
                data["fill_color"] = f"#{fc.rgb}"
    except Exception:
        pass

    # Line color
    try:
        if shape.line and shape.line.color and shape.line.color.rgb:
            data["line_color"] = f"#{shape.line.color.rgb}"
        if shape.line and shape.line.width:
            data["line_width"] = shape.line.width
    except Exception:
        pass

    # Text
    try:
        if shape.has_text_frame:
            texts = []
            for p in shape.text_frame.paragraphs:
                for run in p.runs:
                    texts.append({
                        "text": run.text,
                        "font_size": run.font.size,
                        "font_bold": run.font.bold,
                        "font_color": f"#{run.font.color.rgb}" if run.font.color and run.font.color.rgb else None,
                        "font_name": run.font.name,
                    })
            data["text_runs"] = texts
            data["text"] = shape.text_frame.text
    except Exception:
        pass

    # Chart
    try:
        if shape.has_chart:
            data["has_chart"] = True
            chart = shape.chart
            data["chart_type"] = str(chart.chart_type)
            if chart.has_title:
                data["chart_title"] = chart.chart_title.text_frame.text
    except Exception:
        pass

    return data


def extract_all_shapes(pptx_path: str) -> dict:
    """Extract all shapes from all slides, keyed by name."""
    prs = Presentation(pptx_path)
    shapes_by_slide = {}
    for i, slide in enumerate(prs.slides):
        slide_shapes = {}
        for shape in slide.shapes:
            name = shape.name
            slide_shapes[name] = _extract_shape_data(shape)
        shapes_by_slide[f"slide_{i+1}"] = slide_shapes
    return shapes_by_slide


# ── Diff Computation ──────────────────────────────────────────

def compute_diff(original: dict, modified: dict) -> list[dict]:
    """Compute differences between original and modified PPT shapes.

    Returns list of change dicts:
    {
        "slide": str,
        "shape_name": str,
        "change_type": "color"|"size"|"position"|"text"|"deleted"|"added",
        "property": str,
        "old_value": any,
        "new_value": any,
    }
    """
    changes = []

    all_slides = set(list(original.keys()) + list(modified.keys()))

    for slide_key in sorted(all_slides):
        orig_shapes = original.get(slide_key, {})
        mod_shapes = modified.get(slide_key, {})

        # Deleted shapes
        for name in orig_shapes:
            if name not in mod_shapes:
                changes.append({
                    "slide": slide_key,
                    "shape_name": name,
                    "change_type": "deleted",
                    "property": "shape",
                    "old_value": orig_shapes[name].get("shape_type"),
                    "new_value": None,
                })

        # Added shapes
        for name in mod_shapes:
            if name not in orig_shapes:
                changes.append({
                    "slide": slide_key,
                    "shape_name": name,
                    "change_type": "added",
                    "property": "shape",
                    "old_value": None,
                    "new_value": mod_shapes[name].get("shape_type"),
                })

        # Modified shapes
        for name in orig_shapes:
            if name not in mod_shapes:
                continue
            orig = orig_shapes[name]
            mod = mod_shapes[name]

            # Color changes
            for color_prop in ["fill_color", "line_color"]:
                if orig.get(color_prop) != mod.get(color_prop):
                    if orig.get(color_prop) or mod.get(color_prop):
                        changes.append({
                            "slide": slide_key,
                            "shape_name": name,
                            "change_type": "color",
                            "property": color_prop,
                            "old_value": orig.get(color_prop),
                            "new_value": mod.get(color_prop),
                        })

            # Size changes (>5% threshold)
            for dim in ["width", "height"]:
                ov = orig.get(dim, 0)
                nv = mod.get(dim, 0)
                if ov and nv and abs(nv - ov) / ov > 0.05:
                    changes.append({
                        "slide": slide_key,
                        "shape_name": name,
                        "change_type": "size",
                        "property": dim,
                        "old_value": ov,
                        "new_value": nv,
                        "pct_change": (nv - ov) / ov,
                    })

            # Position changes (>5% of slide width)
            for pos in ["left", "top"]:
                ov = orig.get(pos, 0)
                nv = mod.get(pos, 0)
                if abs(nv - ov) > Emu(Inches(0.2)):
                    changes.append({
                        "slide": slide_key,
                        "shape_name": name,
                        "change_type": "position",
                        "property": pos,
                        "old_value": ov,
                        "new_value": nv,
                    })

            # Text changes
            if orig.get("text") != mod.get("text"):
                if orig.get("text") or mod.get("text"):
                    changes.append({
                        "slide": slide_key,
                        "shape_name": name,
                        "change_type": "text",
                        "property": "text",
                        "old_value": orig.get("text"),
                        "new_value": mod.get("text"),
                    })

            # Font changes
            orig_runs = orig.get("text_runs", [])
            mod_runs = mod.get("text_runs", [])
            for j in range(min(len(orig_runs), len(mod_runs))):
                or_ = orig_runs[j]
                mr_ = mod_runs[j]
                for prop in ["font_size", "font_bold", "font_color", "font_name"]:
                    if or_.get(prop) != mr_.get(prop):
                        changes.append({
                            "slide": slide_key,
                            "shape_name": name,
                            "change_type": "font",
                            "property": prop,
                            "old_value": or_.get(prop),
                            "new_value": mr_.get(prop),
                        })

    return changes


# ── Style Token Update ────────────────────────────────────────

def changes_to_token_updates(changes: list[dict]) -> list[dict]:
    """Convert raw changes to actionable style token updates.

    Maps shape name patterns to style token paths:
      paradigm_epoch_* → ParadigmLayout, ParadigmColor
      paradigm_icon_* → icon system
      *_chart → chart styling
      panel_label_* → PanelLabel
    """
    updates = []

    for change in changes:
        name = change["shape_name"]
        ct = change["change_type"]
        prop = change["property"]

        # Paradigm epoch boxes
        if name.startswith("paradigm_epoch_"):
            if ct == "color" and prop == "fill_color":
                updates.append({
                    "target": "ParadigmColor.SCREEN_BG",
                    "action": "user_prefers_different_fill",
                    "old": change["old_value"],
                    "new": change["new_value"],
                    "confidence": "high",
                })
            elif ct == "size":
                direction = "larger" if change.get("pct_change", 0) > 0 else "smaller"
                updates.append({
                    "target": f"ParadigmLayout.BOX_{prop.upper()}",
                    "action": f"user_wants_{direction}_boxes",
                    "pct_change": change.get("pct_change"),
                    "confidence": "high",
                })

        # Paradigm icons
        elif name.startswith("paradigm_icon_") and not name.endswith("_symbol"):
            if ct == "color":
                updates.append({
                    "target": "ParadigmColor.icon_colors",
                    "action": "user_changed_icon_color",
                    "old": change["old_value"],
                    "new": change["new_value"],
                    "confidence": "medium",
                })
            elif ct == "deleted":
                updates.append({
                    "target": "paradigm_template",
                    "action": "user_removed_icon_strip",
                    "confidence": "high",
                    "note": "User may prefer no icon color strips",
                })

        # Labels
        elif name.startswith("paradigm_label_"):
            if ct == "text":
                updates.append({
                    "target": "content",
                    "action": "user_renamed_epoch",
                    "old": change["old_value"],
                    "new": change["new_value"],
                    "confidence": "content_change",
                })
            elif ct == "font":
                updates.append({
                    "target": f"SlideFont.{prop.upper()}",
                    "action": f"user_changed_{prop}",
                    "old": change["old_value"],
                    "new": change["new_value"],
                    "confidence": "high",
                })

        # Panel labels
        elif name.startswith("panel_label_"):
            if ct == "font":
                updates.append({
                    "target": "PanelLabel",
                    "action": f"user_changed_{prop}",
                    "old": change["old_value"],
                    "new": change["new_value"],
                    "confidence": "high",
                })

        # Charts
        elif "chart" in name:
            if ct == "deleted":
                updates.append({
                    "target": "slide_template",
                    "action": "user_removed_chart",
                    "chart_name": name,
                    "confidence": "high",
                })

        # General: any deletion
        if ct == "deleted":
            updates.append({
                "target": "general",
                "action": "element_removed",
                "shape_name": name,
                "confidence": "medium",
            })

    return updates


# ── Feedback Processing Pipeline ──────────────────────────────

def process_feedback(original_pptx: str, modified_pptx: str,
                     output_dir: Path = None) -> dict:
    """Full feedback pipeline: extract → diff → token updates.

    Returns:
    {
        "timestamp": str,
        "n_changes": int,
        "changes": [dict],
        "token_updates": [dict],
        "summary": str,
    }
    """
    if output_dir is None:
        output_dir = FEEDBACK_DIR
    output_dir.mkdir(parents=True, exist_ok=True)

    # Extract shapes
    orig_shapes = extract_all_shapes(original_pptx)
    mod_shapes = extract_all_shapes(modified_pptx)

    # Compute diff
    changes = compute_diff(orig_shapes, mod_shapes)

    # Generate token updates
    token_updates = changes_to_token_updates(changes)

    # Summary
    change_types = {}
    for c in changes:
        change_types[c["change_type"]] = change_types.get(c["change_type"], 0) + 1

    summary_parts = []
    if change_types.get("color"):
        summary_parts.append(f"{change_types['color']} color changes")
    if change_types.get("size"):
        summary_parts.append(f"{change_types['size']} size changes")
    if change_types.get("position"):
        summary_parts.append(f"{change_types['position']} position changes")
    if change_types.get("text"):
        summary_parts.append(f"{change_types['text']} text edits")
    if change_types.get("deleted"):
        summary_parts.append(f"{change_types['deleted']} elements deleted")
    if change_types.get("added"):
        summary_parts.append(f"{change_types['added']} elements added")
    if change_types.get("font"):
        summary_parts.append(f"{change_types['font']} font changes")

    summary = "; ".join(summary_parts) if summary_parts else "No changes detected"

    result = {
        "timestamp": datetime.now(timezone.utc).isoformat(),
        "original": original_pptx,
        "modified": modified_pptx,
        "n_changes": len(changes),
        "change_summary": change_types,
        "changes": changes,
        "token_updates": token_updates,
        "summary": summary,
    }

    # Save feedback
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    feedback_path = output_dir / f"feedback_{ts}.json"
    with open(feedback_path, "w") as f:
        json.dump(result, f, indent=2, default=str)

    print(f"Feedback saved: {feedback_path}")
    print(f"  Changes: {len(changes)} ({summary})")
    print(f"  Token updates: {len(token_updates)}")

    for update in token_updates[:10]:
        print(f"    → {update['target']}: {update['action']}")

    return result


# ── CLI ───────────────────────────────────────────────────────

if __name__ == "__main__":
    import argparse
    parser = argparse.ArgumentParser(description="Process PPT feedback")
    parser.add_argument("original", help="Original PPTX path")
    parser.add_argument("modified", help="User-modified PPTX path")
    args = parser.parse_args()

    result = process_feedback(args.original, args.modified)
    print(f"\n{result['summary']}")
